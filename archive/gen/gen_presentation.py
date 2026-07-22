"""
Generate the WSmart+ Route results PowerPoint presentation.

Builds a 20-slide deck under assets/windows/ following the agreed structure:

  1.  Cover (title/authors/affiliations, as in the conference abstract)
  2.  Index / agenda (condensed)
  3.  The VRPP problem                          (equation image + variable glossary)
  4.  Objective of this work                    (bullets)
  5.  Routing simulator philosophy              (pipeline diagram + caption)
  6.  Mandatory node selection strategies       (equation + caption)
  7.  Big picture: policy configuration space   (grid diagram + caption)
  8.  Exact methods (BPC + SWC-TCF)             (equation + figure + caption)
  9.  Meta-heuristics & hyper-heuristics        (equation + caption, all algorithms)
  10. CLS vs Fast-TSP route improvers           (equation + figure + caption)
  11. Design of experiments                     (tree diagram + caption)
  12. Pareto front plot (log X, per-region styled fronts, figure caption)
  13. Strategy trade-off bubble chart (figure caption)
  14. Per-scenario heatmaps (30 days, figure caption)
  15. Overflow + kg/km policy × scenario heatmaps (90 days, side legend + caption)
  16. Route improver bubble chart (figure caption)
  17. Full results table (user-selected horizon, or all horizons — table caption)
  18. Conclusions, limitations & future work    (radar figure)
  19. Acknowledgements
  20. End / Q&A (figure)

Slide text and captions live in json/presentation_content.json; equations
("equation" key) are rendered as transparent-background images via
matplotlib mathtext (render_equation_image) rather than native OOXML math
objects — see the module docstring note below render_equation_image for why.
Result figures are pulled from the simulation analysis figures directories
(30-day set by default; individual figure slides may override via
"figures_dir"). The full results table (slide 17) is built directly from the
horizon CSV(s) referenced in json/simulation_analysis_config.json.

A per-slide speaker script can also be generated as a .docx (see
gen_speaker_script / --speaker-script), rendered via docxtpl from a template
under archive/gen/templates/.

Usage
-----
    uv run python archive/gen/gen_presentation.py
    uv run python archive/gen/gen_presentation.py \\
        --figures-dir public/figures/simulation/30d \\
        --out assets/windows/wsmart_route_results.pptx \\
        --author "Afonso Fernandes" \\
        --coauthors "Jane Doe;John Smith" \\
        --groups "ISR Coimbra;INESC-ID Lisboa" \\
        --results-table 30d \\
        --speaker-script
"""

from __future__ import annotations

import argparse
import math
import re
import shutil
import tempfile
import textwrap
import xml.etree.ElementTree as ET
from pathlib import Path

import matplotlib
import numpy as np

matplotlib.use("Agg")
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
from gen_simulation_analysis import (  # pyrefly: ignore [missing-import]
    KGKM_LABEL,
    _group_spans,
    aggregate,
    build_context,
    build_full_results_matrix,
    disp,
    filter_data,
    load_horizon_csv,
)
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationClass
from pptx.util import Emu, Inches, Pt
from report_utils import load_json  # pyrefly: ignore [missing-import]

# ── Equations as images (matplotlib mathtext) ───────────────────────────────
# An earlier version embedded equations as native, editable OOXML <m:oMath>
# objects (pandoc LaTeX -> OMML, wrapped in mc:AlternateContent/a14:m). That
# renders and edits correctly in desktop PowerPoint, but both LibreOffice
# Impress and the PowerPoint web app either blank the math branch or refuse to
# open it for editing — i.e. it fails for exactly the viewers this deck also
# needs to work in. A plain raster image is strictly more portable: every
# viewer just displays a picture. matplotlib's built-in "mathtext" renders the
# large majority of the LaTeX used in this deck (\mathcal, \mathbb, \bar,
# \sum, \min, \exists, \Rightarrow, \geq, sub/superscripts, ...) without
# requiring a system LaTeX install, and — unlike real LaTeX — tolerates a
# label written in \mathbf{...} that contains literal spaces/punctuation, so
# a bold step-label can be authored directly inside the math markup.
#
# A line may mix plain text and math freely by wrapping only the math spans in
# `$...$` (matplotlib supports multiple math spans per string); anything
# outside `$...$` renders as plain (non-italic) text. A line with no `$` at
# all is treated as a pure formula and auto-wrapped. Whether a line renders
# bold is inferred from its shape: a line that opens with `$` is a bare
# formula (kept at regular weight, matching the historical VRPP equation
# style); a line that starts with plain text (e.g. a "Step N — ..." label) is
# bolded in full, which reads naturally since the label dominates the line.
def _prepare_equation_line(raw: str) -> tuple[str, bool]:
    line = raw.replace(r"\textbf{", r"\mathbf{")  # legacy authoring safety net
    if "$" not in line:
        return f"${line}$", False
    return line, not line.lstrip().startswith("$")


def render_equation_image(
    lines: list[str],
    out_path: Path,
    width_in: float,
    fontsize: float = 20,
    line_h_in: float = 0.62,
    color: str = "#1F2D3D",
    align: str = "left",
) -> Path:
    """Render equation `lines` (mixed plain text + `$...$` mathtext) to a
    transparent PNG sized at exactly `width_in` x (`line_h_in` * len(lines))
    inches, so embedding it at that same size in the slide reproduces
    `fontsize` as literal points — matching how the OMML equations used to be
    sized via `size_pt`.
    """
    n = max(len(lines), 1)
    height_in = line_h_in * n
    fig = plt.figure(figsize=(width_in, height_in), dpi=220)
    fig.patch.set_alpha(0.0)
    x = {"left": 0.012, "center": 0.5, "right": 0.988}[align]
    for i, raw in enumerate(lines):
        text, bold = _prepare_equation_line(raw)
        y = 1.0 - (i + 0.5) / n
        fig.text(
            x, y, text, ha=align, va="center", fontsize=fontsize, color=color,
            fontweight="bold" if bold else "normal",
        )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=220, transparent=True)
    plt.close(fig)
    return out_path


def render_glossary_image(
    items: list[tuple[str, str]],
    out_path: Path,
    width_in: float,
    ncols: int = 2,
    fontsize: float = 12.5,
    color: str = "#1F2D3D",
    max_lines: int = 2,
) -> tuple[Path, float]:
    """Render a "$symbol$ — meaning" variable glossary as a transparent PNG,
    laid out in `ncols` columns of left-aligned rows spanning `width_in`
    inches (see `render_equation_image` for why images, not native equation
    objects). The `meaning` half of each entry is word-wrapped (up to
    `max_lines`) to the column width so it cannot run into the next column —
    matplotlib's mathtext does not wrap on its own.

    Unlike `render_equation_image` (fixed line count), a glossary's row count
    that actually needs 2 wrapped lines varies with content, so the image
    height is *derived* from the wrapped layout (row-by-row, each row sized to
    its tallest cell) rather than a height the caller must pre-guess — a
    pre-guessed, too-short height silently overlapped adjacent rows.

    Returns the output path and the resulting image height in inches.
    """
    nrows = math.ceil(len(items) / ncols)
    col_w = 1.0 / ncols
    col_w_in = width_in * col_w
    # Crude average-glyph-width estimate (~0.5 * point size, in inches) to size
    # the wrap width — good enough since mismatches only cost a slightly early
    # or late line break, never overflow into the next column.
    char_w_in = fontsize * 0.0086
    max_chars = max(int((col_w_in - 0.1) / char_w_in), 10)
    wrapped_items = []
    for symbol, meaning in items:
        wrapped = (textwrap.wrap(meaning, width=max_chars) or [meaning])[:max_lines]
        wrapped_items.append((symbol, wrapped))

    line_unit_in = fontsize * 0.034  # matches the equation renderer's pt->inch line pitch
    row_lines = [1] * nrows
    for i, (_, wrapped) in enumerate(wrapped_items):
        _, row = divmod(i, nrows)
        row_lines[row] = max(row_lines[row], len(wrapped))
    row_h_in = [n * line_unit_in for n in row_lines]
    height_in = sum(row_h_in)
    row_top_in = [sum(row_h_in[:r]) for r in range(nrows)]

    fig = plt.figure(figsize=(width_in, height_in), dpi=220)
    fig.patch.set_alpha(0.0)
    for i, (symbol, wrapped) in enumerate(wrapped_items):
        col, row = divmod(i, nrows)
        x = col * col_w + 0.008
        y = 1.0 - (row_top_in[row] + row_h_in[row] / 2) / height_in
        text = f"{symbol} — {wrapped[0]}" + "".join(f"\n{extra}" for extra in wrapped[1:])
        fig.text(x, y, text, ha="left", va="center", fontsize=fontsize, color=color, linespacing=1.2)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=220, transparent=True)
    plt.close(fig)
    return out_path, height_in


# 16:9 slide geometry
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT = RGBColor(0x2E, 0x74, 0xB5)
DARK = RGBColor(0x1F, 0x2D, 0x3D)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
LIGHT_TXT = RGBColor(0xC9, 0xD6, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PIPELINE_STAGES = [
    ("Mandatory bin\nselection strategy", "is there a collection today?"),
    ("Route\nconstructor", "build the routes"),
    ("Acceptance criterion\n(optional)", "constructor-dependent"),
    ("Route improver\n(optional)", "post-optimisation"),
]

ASSETS_DIR = Path(__file__).resolve().parent.parent.parent / "assets" / "images"


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _add_connector(slide, x, y, cx, cy, color="2E74B5", flip_v=False, flip_h=False, width_pt=1.5) -> None:
    """Inject a straight arrow connector (cxnSp) directly into the slide spTree."""
    PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    AML = "http://schemas.openxmlformats.org/drawingml/2006/main"
    new_id = len(slide.shapes) + 900
    flip_attr = (' flipH="1"' if flip_h else "") + (' flipV="1"' if flip_v else "")
    w_emu = int(width_pt * 12700)
    xml = (
        f'<p:cxnSp xmlns:p="{PML}" xmlns:a="{AML}">'
        f'<p:nvCxnSpPr><p:cNvPr id="{new_id}" name="Connector {new_id}"/>'
        f'<p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm{flip_attr}>'
        f'<a:off x="{int(x)}" y="{int(y)}"/><a:ext cx="{int(cx)}" cy="{int(cy)}"/>'
        f'</a:xfrm><a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{w_emu}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'<a:tailEnd type="arrow"/></a:ln></p:spPr></p:cxnSp>'
    )
    slide.shapes._spTree.append(etree.fromstring(xml))


def _wrap_label(label, width: int) -> str:
    """Wrap a table header/cell label onto multiple lines so it fits its cell."""
    return "\n".join(textwrap.wrap(str(label), width=width)) or str(label)


_CELL_RE = re.compile(r"([-\d.]+(?:±[-\d.]+)?)\s*ov\s*(?:<br>|\n)\s*([-\d.]+(?:±[-\d.]+)?)\s*kg/km")


def _parse_result_cell(text: str):
    """Split a formatted result cell ('X ov<br>Y kg/km') into (overflow, kgkm) as (raw_str, numeric_mean)."""
    m = _CELL_RE.match(text)
    if not m:
        return None, None

    def _pair(raw: str):
        try:
            return raw, float(raw.split("±")[0])
        except ValueError:
            return raw, None

    return _pair(m.group(1)), _pair(m.group(2))


def compute_global_best(row_keys: list[tuple], col_keys: list[tuple], cells: dict[tuple, str]) -> dict:
    """Best (lowest overflow / highest kg-km) column per row, across the FULL column set.

    Computed once over every strategy/constructor/improver combination so that,
    when the table is later split into per-strategy partial images, the same
    single best cell is highlighted rather than one (possibly different) best
    per partition.
    """
    best: dict = {}
    for rk in row_keys:
        ov_vals, kg_vals = {}, {}
        for ck in col_keys:
            ov, kg = _parse_result_cell(cells.get((rk, ck), "—"))
            if ov and ov[1] is not None:
                ov_vals[ck] = ov[1]
            if kg and kg[1] is not None:
                kg_vals[ck] = kg[1]
        best[rk] = {
            "ov": min(ov_vals, key=ov_vals.get) if ov_vals else None,  # pyrefly: ignore [no-matching-overload]
            "kg": max(kg_vals, key=kg_vals.get) if kg_vals else None,  # pyrefly: ignore [no-matching-overload]
        }
    return best


def render_hier_table_image(
    row_keys: list[tuple],
    col_keys: list[tuple],
    cells: dict[tuple, str],
    row_labels: list[str],
    out_path: Path,
    col_lookup_keys: list[tuple] | None = None,
    partition_label: str | None = None,
    target_size_in: tuple[float, float] | None = None,
    global_best: dict | None = None,
    corner_note: str | None = None,
) -> Path:
    """
    Render a hierarchical (merged-header) results table as a PNG, in the style
    of a LaTeX multi-row/multi-column table: row groups (region > graph size >
    distribution) and column groups (strategy > constructor > improver, plus
    horizon if present) are drawn as merged header spans.

    `col_lookup_keys`, if given, is parallel to `col_keys` and used instead of
    it to look cells up in `cells` — lets a hierarchy level be dropped from the
    displayed header (e.g. when the table has been partitioned by that level)
    while the underlying data lookup still uses the full key.

    `target_size_in`, if given, stretches the table to exactly fill that
    (width, height) in inches instead of using its own natural aspect ratio —
    lets a narrow partial table still occupy the full slide width.

    `global_best`, if given, overrides the per-partition best-cell computation
    with a row_key -> {"ov": lookup_key, "kg": lookup_key} mapping computed
    once across every partition (see `compute_global_best`).
    """
    col_lookup_keys = col_lookup_keys if col_lookup_keys is not None else col_keys
    n_row_levels = len(row_labels)
    n_col_levels = len(col_keys[0])
    n_rows = len(row_keys)
    n_cols = len(col_keys)

    # Reserve an extra banner row at the top when a partition label is supplied
    banner_h = 0.5 if partition_label else 0.0
    header_h = 0.5
    # Per-level row-label widths: level 0 (region) needs room for names like "Figueira
    # da Foz", level 1 (graph size) just a short number, level 2 (distribution) a medium
    # word. Weighted rather than equal-width, and always a FRACTION of the target width
    # (not fixed inches) so they don't eat a growing share of a shrunk target box.
    level_weights = ([1.6, 0.7, 1.1] + [1.0] * n_row_levels)[:n_row_levels]
    weight_sum = sum(level_weights)

    if target_size_in:
        fig_w, fig_h = target_size_in
        # Headers/banner must not swallow the whole target height when many rows are
        # stacked into a small box (e.g. a per-strategy partial table) — cap their
        # share of fig_h so every data row keeps a visible sliver of height.
        if partition_label:
            banner_h = min(banner_h, fig_h * 0.16)
        header_h = min(header_h, fig_h * 0.14)
        x0 = fig_w * 0.16
        label_ws = [x0 * w / weight_sum for w in level_weights]
        y0 = banner_h + header_h * n_col_levels
        cell_w = (fig_w - x0) / max(n_cols, 1)
        cell_h = (fig_h - y0) / max(n_rows, 1)
        total_w, total_h = fig_w, fig_h
    else:
        cell_w, cell_h = 1.05, 0.5
        x0 = 3.6
        label_ws = [x0 * w / weight_sum for w in level_weights]
        y0 = banner_h + header_h * n_col_levels
        total_h = banner_h + header_h * n_col_levels + cell_h * n_rows
        total_w = x0 + cell_w * n_cols
        fig_w = min(total_w, 42)
        fig_h = min(total_h, 32)
    label_lefts = [sum(label_ws[:i]) for i in range(n_row_levels)]

    fontsize = max(5.5, min(11, 300 / max(n_cols, 1)))
    # Two stacked values (overflows over kg/km) must fit the cell: each line
    # gets ~cell_h/2 of height and the full cell width for one value string.
    fontsize = min(fontsize, cell_h * 72 * 0.30, cell_w * 72 * 0.185)

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    ax.set_xlim(0, total_w)
    ax.set_ylim(0, total_h)
    ax.invert_yaxis()
    ax.axis("off")

    # Banner row spanning all columns (including the row-label area)
    if partition_label:
        ax.add_patch(mpatches.Rectangle((0, 0), total_w, banner_h, facecolor="#0D1B2A", edgecolor="none"))
        ax.text(
            total_w / 2, banner_h / 2, partition_label, ha="center", va="center",
            fontsize=fontsize + 2, color="white", fontweight="bold",
        )

    # Corner note (top-left, above the row labels / left of the column headers):
    # explains that the top value in every cell is overflows, bottom is kg/km.
    if corner_note:
        ax.text(
            x0 / 2, banner_h + (header_h * n_col_levels) / 2, corner_note, ha="center", va="center",
            fontsize=fontsize * 0.85, color="#1F2D3D", fontweight="bold", linespacing=1.6,
        )

    header_colors = ["#1F2D3D", "#2E74B5", "#5A6A7A", "#8A9BB0"]
    for lvl in range(n_col_levels):
        y_top = banner_h + lvl * header_h
        for start, end, label in _group_spans(col_keys, lvl):
            xs, xe = x0 + start * cell_w, x0 + end * cell_w
            ax.add_patch(
                mpatches.Rectangle(
                    (xs, y_top), xe - xs, header_h,
                    facecolor=header_colors[lvl % len(header_colors)], edgecolor="white", linewidth=0.8,
                )
            )
            wrap_chars = max(6, int((xe - xs) / cell_w) * 10)
            ax.text(
                (xs + xe) / 2, y_top + header_h / 2, _wrap_label(disp(label), wrap_chars),
                ha="center", va="center", fontsize=fontsize, color="white", fontweight="bold",
            )

    for lvl in range(n_row_levels):
        x_left, lw = label_lefts[lvl], label_ws[lvl]
        # Single-line label sized to fit the cell width, so it can never look like it
        # straddles two rows (a wrapped 2-line label centred on a 2-row span can land
        # one line per row and read as two separate cells).
        lvl_fontsize = min(fontsize, lw * 72 * 0.16)
        for start, end, label in _group_spans(row_keys, lvl):
            ys, ye = y0 + start * cell_h, y0 + end * cell_h
            ax.add_patch(
                mpatches.Rectangle((x_left, ys), lw, ye - ys, facecolor="#F0F4FA", edgecolor="#5A6A7A", linewidth=0.6)
            )
            ax.text(
                x_left + lw / 2, (ys + ye) / 2, str(label),
                ha="center", va="center", fontsize=lvl_fontsize, color="#1F2D3D", fontweight="bold",
            )
            # Separator across the data area (not the row-label area) at every
            # innermost row-group boundary (e.g. between the Empirical and Gamma-3 rows).
            if lvl == n_row_levels - 1 and start > 0:
                ax.plot([x0, total_w], [ys, ys], color="#5A6A7A", linewidth=1.2, zorder=5)

    for ri, rk in enumerate(row_keys):
        best_ov_key = None
        best_kg_key = None
        best_ov_ci = None
        best_kg_ci = None
        if global_best is not None:
            best_ov_key = global_best.get(rk, {}).get("ov")
            best_kg_key = global_best.get(rk, {}).get("kg")
            parsed = {ci: _parse_result_cell(cells.get((rk, lk), "—")) for ci, lk in enumerate(col_lookup_keys)}
        else:
            parsed = {ci: _parse_result_cell(cells.get((rk, lk), "—")) for ci, lk in enumerate(col_lookup_keys)}
            ov_vals = {ci: v[0][1] for ci, v in parsed.items() if v[0] and v[0][1] is not None}
            kg_vals = {ci: v[1][1] for ci, v in parsed.items() if v[1] and v[1][1] is not None}
            best_ov_ci = min(ov_vals, key=ov_vals.get) if ov_vals else None  # pyrefly: ignore [no-matching-overload]
            best_kg_ci = max(kg_vals, key=kg_vals.get) if kg_vals else None  # pyrefly: ignore [no-matching-overload]
        for ci in range(n_cols):
            xs, ys = x0 + ci * cell_w, y0 + ri * cell_h
            ax.add_patch(mpatches.Rectangle((xs, ys), cell_w, cell_h, fill=False, edgecolor="#CCCCCC", linewidth=0.4))
            ov, kg = parsed.get(ci, (None, None))
            if ov is None and kg is None:
                ax.text(xs + cell_w / 2, ys + cell_h / 2, "—", ha="center", va="center", fontsize=fontsize * 0.85)
                continue
            # Dotted separator between the overflow (top) and kg/km (bottom) values in the cell.
            ax.plot(
                [xs + cell_w * 0.1, xs + cell_w * 0.9], [ys + cell_h * 0.5, ys + cell_h * 0.5],
                linestyle=":", color="#AAAAAA", linewidth=0.8,
            )
            if global_best is not None:
                is_best_ov = best_ov_key is not None and col_lookup_keys[ci] == best_ov_key
                is_best_kg = best_kg_key is not None and col_lookup_keys[ci] == best_kg_key
            else:
                is_best_ov = ci == best_ov_ci
                is_best_kg = ci == best_kg_ci
            # Best half-cells get a light-green background behind the bold green value.
            if is_best_ov:
                ax.add_patch(mpatches.Rectangle(
                    (xs, ys), cell_w, cell_h * 0.5, facecolor="#D5EEDC", edgecolor="none", zorder=0.6,
                ))
            if is_best_kg:
                ax.add_patch(mpatches.Rectangle(
                    (xs, ys + cell_h * 0.5), cell_w, cell_h * 0.5, facecolor="#D5EEDC", edgecolor="none", zorder=0.6,
                ))
            ov_color, ov_weight = ("#1A7A34", "bold") if is_best_ov else ("#333333", "bold")
            kg_color, kg_weight = ("#1A7A34", "bold") if is_best_kg else ("#333333", "bold")
            ax.text(
                xs + cell_w / 2, ys + cell_h * 0.27, ov[0] if ov else "—", ha="center", va="center",
                fontsize=fontsize, color=ov_color, fontweight=ov_weight,
            )
            ax.text(
                xs + cell_w / 2, ys + cell_h * 0.74, kg[0] if kg else "—", ha="center", va="center",
                fontsize=fontsize, color=kg_color, fontweight=kg_weight,
            )

    fig.savefig(out_path, dpi=180, bbox_inches=None, pad_inches=0, facecolor="white")
    plt.close(fig)
    return out_path


TEMPLATES_DIR = Path(__file__).resolve().parent / "templates"


def gen_speaker_script(deck_title: str, author: str, slides: list[dict], out_path: Path) -> Path:
    """Render a per-slide speaker script .docx from the docxtpl template."""
    from docxtpl import DocxTemplate

    tpl = DocxTemplate(str(TEMPLATES_DIR / "speaker_script.docx"))
    tpl.render({"deck_title": deck_title, "author": author, "slides": slides})
    out_path.parent.mkdir(parents=True, exist_ok=True)
    tpl.save(str(out_path))
    return out_path


_INLINE_CODE_RE = re.compile(r"`([^`]+)`")
_INLINE_BOLD_RE = re.compile(r"\*\*([^*]+)\*\*")


def _add_markdown_paragraph(doc, text: str, *, size_pt: int = 11) -> None:
    """Add one paragraph, splitting `**bold**` and `` `code` `` spans into separate runs
    (a small hand-rolled inline-markdown renderer — good enough for the appendix's prose,
    not a general Markdown engine)."""
    from docx.shared import Pt as DocxPt

    p = doc.add_paragraph()
    p.paragraph_format.space_after = DocxPt(8)
    # Tokenize on bold/code spans, longest-match-first isn't needed since the two never nest here.
    pos = 0
    tokens: list[tuple[str, str]] = []  # (kind, text): kind in {"plain","bold","code"}
    combined_re = re.compile(r"\*\*([^*]+)\*\*|`([^`]+)`")
    for m in combined_re.finditer(text):
        if m.start() > pos:
            tokens.append(("plain", text[pos : m.start()]))
        if m.group(1) is not None:
            tokens.append(("bold", m.group(1)))
        else:
            tokens.append(("code", m.group(2)))
        pos = m.end()
    if pos < len(text):
        tokens.append(("plain", text[pos:]))
    if not tokens:
        tokens = [("plain", text)]
    for kind, chunk in tokens:
        r = p.add_run(chunk)
        r.font.size = DocxPt(size_pt)
        if kind == "bold":
            r.font.bold = True
        elif kind == "code":
            r.font.name = "Consolas"
            r.font.size = DocxPt(size_pt - 0.5)


def append_markdown_appendix(docx_path: Path, md_path: Path, title: str = "Appendix — Q&A Preparation Notes") -> None:
    """Append a Markdown file to an existing speaker-script .docx as a new titled section.

    Supports a small subset of Markdown sufficient for the appendix content: `#`/`##`/`###`
    headings (mapped to Word Heading 1/2/3, offset so the appendix's own top-level heading
    is always Heading 1), blank-line-separated paragraphs, `**bold**` and `` `code` `` inline
    spans, and `- `/`* ` bullet list items. Not a general CommonMark implementation.
    """
    from docx import Document
    from docx.shared import Pt as DocxPt

    doc = Document(str(docx_path))
    doc.add_page_break()
    doc.add_heading(title, level=1)

    md_text = md_path.read_text(encoding="utf-8")
    # Split into blocks on blank lines, keeping heading lines as their own block.
    lines = md_text.splitlines()
    para_buf: list[str] = []

    def _flush() -> None:
        if para_buf:
            text = " ".join(line.strip() for line in para_buf if line.strip())
            if text:
                _add_markdown_paragraph(doc, text)
            para_buf.clear()

    for raw_line in lines:
        line = raw_line.rstrip()
        heading_m = re.match(r"^(#{1,4})\s+(.*)$", line)
        bullet_m = re.match(r"^[-*]\s+(.*)$", line)
        if not line.strip():
            _flush()
        elif heading_m:
            _flush()
            level = min(len(heading_m.group(1)) + 1, 4)  # md "#" (h1) -> docx Heading 2, etc.
            doc.add_heading(heading_m.group(2).strip(), level=level)
        elif bullet_m:
            _flush()
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = DocxPt(4)
            for kind, chunk in _split_inline_markdown(bullet_m.group(1)):
                r = p.add_run(chunk)
                r.font.size = DocxPt(11)
                if kind == "bold":
                    r.font.bold = True
                elif kind == "code":
                    r.font.name = "Consolas"
                    r.font.size = DocxPt(10.5)
        else:
            para_buf.append(line)
    _flush()
    doc.save(str(docx_path))


def _split_inline_markdown(text: str) -> list[tuple[str, str]]:
    """Tokenize `**bold**` / `` `code` `` spans out of a line; shared by bullet + paragraph rendering."""
    tokens: list[tuple[str, str]] = []
    pos = 0
    combined_re = re.compile(r"\*\*([^*]+)\*\*|`([^`]+)`")
    for m in combined_re.finditer(text):
        if m.start() > pos:
            tokens.append(("plain", text[pos : m.start()]))
        tokens.append(("bold", m.group(1)) if m.group(1) is not None else ("code", m.group(2)))
        pos = m.end()
    if pos < len(text):
        tokens.append(("plain", text[pos:]))
    return tokens or [("plain", text)]


class DeckBuilder:
    def __init__(
        self,
        content: dict,
        figures_dir: Path,
        author: str | None,
        coauthors: list[str] | None = None,
        groups: list[str] | None = None,
        results_table: str = "30d",
        results_table_split: str = "none",
    ):
        self.content = content
        self.figures_dir = figures_dir
        self.author = author or content["author"]
        self.coauthors = coauthors if coauthors is not None else content.get("coauthors", [])
        self.groups = groups if groups is not None else content.get("research_groups", [])
        self.results_table = results_table
        self.results_table_split = results_table_split
        self._results_config: dict = {}
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]
        self._tmp = Path(tempfile.mkdtemp(prefix="wsr_pptx_"))
        self._eq_count = 0
        self._fig_count = 0
        self._tab_count = 0
        self.slide_scripts: list[dict] = []

    def _record_script(self, title: str, paragraphs: list[str]) -> None:
        """Append a speaker-script entry for the slide just built."""
        self.slide_scripts.append(
            {
                "number": len(self.slide_scripts) + 1,
                "title": title,
                "script": "\n\n".join(p for p in paragraphs if p),
            }
        )

    # ── Building blocks ─────────────────────────────────────────────────────────

    def _new_slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _title_bar(self, slide, title: str) -> None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))
        _fill(bar, DARK)
        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE

    def _bullets(self, slide, bullets: list[str], left=None, top=None, width=None, height=None, size=16, gap=10) -> None:
        if not bullets:
            return
        left = left if left is not None else Inches(0.6)
        top = top if top is not None else Inches(1.3)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        _, tf = _textbox(slide, left, top, width, height or (SLIDE_H - top - Inches(0.4)))
        for i, text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {text}"
            p.font.size = Pt(size)
            p.font.color.rgb = DARK
            p.space_after = Pt(gap)

    def _picture_fit(self, slide, path: Path, left, top, max_w, max_h) -> None:
        from PIL import Image

        with Image.open(path) as im:
            w_px, h_px = im.size
        ratio = min(max_w / w_px, max_h / h_px)
        w, h = int(w_px * ratio), int(h_px * ratio)
        slide.shapes.add_picture(str(path), left + Emu(int((max_w - w) / 2)), top + Emu(int((max_h - h) / 2)), w, h)

    def _caption_box(self, slide, label: str, text: str, top=None, left=None, width=None, height=None) -> None:
        """A numbered caption ('**Figure N:** ...' / '**Equation N:** ...' / '**Table N:** ...')."""
        top = top if top is not None else SLIDE_H - Inches(0.55)
        left = left if left is not None else Inches(0.6)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        box_h = height if height is not None else Inches(0.9)
        _, tf = _textbox(slide, left, top, width, box_h)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{label}: "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = DARK
        r2 = p.add_run()
        r2.text = text
        r2.font.italic = True
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED

    def _figure_caption(self, slide, text: str, top=None, left=None, width=None, height=None) -> None:
        self._fig_count += 1
        self._caption_box(slide, f"Figure {self._fig_count}", text, top=top, left=left, width=width, height=height)

    def _equation_caption(self, slide, text: str, top=None, left=None, width=None) -> None:
        self._caption_box(slide, f"Equation {self._eq_count}", text, top=top, left=left, width=width)

    def _table_caption(self, slide, text: str, top=None) -> None:
        self._tab_count += 1
        self._caption_box(slide, f"Table {self._tab_count}", text, top=top)

    def _equation_focus(self, slide, lines: list[str], left=None, top=None, width=None, size_pt: int = 22, line_h: float = 0.62):
        """Place the equation image (render_equation_image) as the slide's visual focus, inside a light rounded panel."""
        self._eq_count += 1
        left = left if left is not None else Inches(0.6)
        top = top if top is not None else Inches(1.25)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        area_h = Inches(line_h * len(lines) + 0.4)
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, area_h)
        _fill(band, RGBColor(0xF0, 0xF4, 0xFA))
        band.shadow.inherit = False
        inner_w = width - Inches(0.4)
        img_h = Inches(line_h * len(lines))
        img_path = render_equation_image(
            lines, self._tmp / f"eq_{self._eq_count}.png",
            width_in=Emu(inner_w).inches, fontsize=size_pt, line_h_in=line_h,
        )
        img_top = top + int((area_h - img_h) / 2)
        slide.shapes.add_picture(str(img_path), left + Inches(0.2), img_top, inner_w, img_h)
        return top + area_h + Inches(0.2)

    def _variable_glossary(
        self, slide, items: list[tuple[str, str]], left=None, top=None, width=None, max_height=None, ncols: int = 2,
        fontsize: float = 12.5,
    ):
        """Place a "$symbol$ — meaning" variable glossary below an equation (render_glossary_image).
        The glossary image is sized to fit its wrapped content exactly (see render_glossary_image);
        `max_height`, if given, uniformly shrinks it (rare) rather than letting it overflow the slide.
        """
        if not items:
            return top
        left = left if left is not None else Inches(0.6)
        top = top if top is not None else Inches(1.25)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        header_h = Inches(0.32)
        head_box = slide.shapes.add_textbox(left, top, width, header_h)
        hp = head_box.text_frame.paragraphs[0]
        hp.text = "Variables"
        hp.font.size = Pt(13)
        hp.font.bold = True
        hp.font.color.rgb = MUTED
        img_top = top + header_h
        img_path, img_h_in = render_glossary_image(
            items, self._tmp / f"glossary_{self._eq_count}.png",
            width_in=Emu(width).inches, ncols=ncols, fontsize=fontsize,
        )
        img_w = width
        img_h = Inches(img_h_in)
        budget = max_height - header_h if max_height is not None else None
        if budget is not None and img_h > budget:
            scale = budget / img_h
            img_h = budget
            img_w = int(width * scale)
        slide.shapes.add_picture(str(img_path), left, img_top, img_w, img_h)
        return img_top + img_h + Inches(0.15)

    def _pipeline_diagram(self, slide):
        """Draw the policy pipeline as chevron stages with sub-labels. Returns the diagram's bottom y."""
        n = len(PIPELINE_STAGES)
        gap = Inches(0.25)
        left0 = Inches(0.5)
        total_w = SLIDE_W - Inches(1.0)
        stage_w = int((total_w - gap * (n - 1)) / n)
        top = Inches(1.9)
        h = Inches(1.5)
        sub_h = Inches(1.1)
        for i, (name, caption) in enumerate(PIPELINE_STAGES):
            left = left0 + i * (stage_w + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, stage_w, h)  # pyrefly: ignore [bad-argument-type]
            optional = "optional" in name
            _fill(shape, ACCENT if not optional else MUTED)
            shape.shadow.inherit = False
            tf = shape.text_frame
            tf.word_wrap = True
            for li, line in enumerate(name.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            _, ctf = _textbox(slide, left, top + h + Inches(0.15), stage_w, sub_h)
            for li, line in enumerate(caption.split("\n")):
                p = ctf.paragraphs[0] if li == 0 else ctf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.color.rgb = MUTED
                p.alignment = PP_ALIGN.CENTER
        return top + h + Inches(0.15) + sub_h

    def _policy_grid_diagram(self, slide) -> None:
        """Draw the Mandatory Selection x Route Constructor x Route Improver big picture."""
        columns = [
            (
                "Mandatory Selection",
                [
                    "Look-Ahead (LA)",
                    "Last-Minute (LM, CF70)",
                    "Last-Minute (LM, CF90)",
                    "Service-Level (SL1)",
                    "Service-Level (SL2)",
                ],
                "flat",
            ),
            (
                "Route Constructor",
                [
                    ("Exact Methods", ["BPC", "SWC-TCF"]),
                    ("Meta-Heuristics", ["HGS", "ALNS", "SANS", "PG-CLNS", "PSOMA"]),
                    ("Hyper-Heuristics", ["ACO-HH"]),
                ],
                "grouped",
            ),
            ("Route Improver", ["Fast-TSP", "Classical Local Search (CLS)"], "flat"),
        ]
        top, bottom = Inches(1.15), Inches(6.5)
        gap = Inches(0.3)
        left0 = Inches(0.6)
        col_w = int((SLIDE_W - Inches(1.2) - gap * (len(columns) - 1)) / len(columns))
        header_h = Inches(0.55)
        for ci, (header, items, kind) in enumerate(columns):
            left = left0 + ci * (col_w + gap)
            hd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, header_h)
            _fill(hd, DARK)
            hd.shadow.inherit = False
            p = hd.text_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            body_top = top + header_h + Inches(0.15)
            body_h = bottom - body_top
            if kind == "flat":
                item_gap = Inches(0.12)
                item_h = int((body_h - item_gap * (len(items) - 1)) / len(items))
                for ii, label in enumerate(items):
                    it_top = body_top + ii * (item_h + item_gap)
                    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, it_top, col_w, item_h)
                    _fill(box, ACCENT)
                    box.shadow.inherit = False
                    p = box.text_frame.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(13)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
            else:
                grp_gap = Inches(0.15)
                grp_h = int((body_h - grp_gap * (len(items) - 1)) / len(items))
                for gi, (grp_label, sub_items) in enumerate(items):
                    grp_top = body_top + gi * (grp_h + grp_gap)
                    grp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, grp_top, col_w, grp_h)
                    grp.fill.background()
                    grp.line.color.rgb = MUTED
                    grp.line.width = Pt(1)
                    grp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                    grp.shadow.inherit = False
                    tf = grp.text_frame
                    tf.word_wrap = True
                    tf.margin_top = Inches(0.05)
                    p = tf.paragraphs[0]
                    p.text = grp_label
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = ACCENT
                    p.alignment = PP_ALIGN.CENTER
                    pi = tf.add_paragraph()
                    pi.text = "  ·  ".join(sub_items)
                    pi.font.size = Pt(11)
                    pi.font.color.rgb = DARK
                    pi.alignment = PP_ALIGN.CENTER

    def _algo_taxonomy_diagram(self, slide) -> None:
        """Draw the Exact / Meta-Heuristic / Hyper-Heuristic route-constructor taxonomy."""
        groups = [
            ("Exact Methods", ACCENT, [
                "Branch-Price-and-Cut\n(BPC)",
                "Smart Waste Collection\nTwo-Commodity Flow (SWC-TCF)",
            ]),
            ("Meta-Heuristics", RGBColor(0x3E, 0x8E, 0x41), [
                "Hybrid Genetic Search (HGS)",
                "Adaptive Large Neighborhood Search (ALNS)",
                "Simulated Annealing Neighborhood Search (SANS)",
                "Policy-Gradient Cooperative LNS (PG-CLNS)",
                "Particle Swarm Optimisation Memetic Algorithm (PSOMA)",
            ]),
            ("Hyper-Heuristics", RGBColor(0xB0, 0x6A, 0x2E), [
                "Ant Colony Optimisation Hyper-Heuristic (ACO-HH)",
            ]),
        ]
        top, bottom = Inches(1.25), Inches(4.15)
        gap = Inches(0.35)
        left0 = Inches(0.6)
        col_w = int((SLIDE_W - Inches(1.2) - gap * (len(groups) - 1)) / len(groups))
        header_h = Inches(0.6)
        for ci, (header, color, items) in enumerate(groups):
            left = left0 + ci * (col_w + gap)
            hd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, header_h)
            _fill(hd, color)
            hd.shadow.inherit = False
            p = hd.text_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            body_top = top + header_h + Inches(0.15)
            body_h = bottom - body_top
            item_gap = Inches(0.12)
            item_h = int((body_h - item_gap * (len(items) - 1)) / len(items))
            for ii, label in enumerate(items):
                it_top = body_top + ii * (item_h + item_gap)
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, it_top, col_w, item_h)
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA)
                box.line.color.rgb = color
                box.line.width = Pt(1.25)
                box.shadow.inherit = False
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER

    def _doe_tree_diagram(self, slide) -> int:
        """Draw the design-of-experiments tree: horizon -> scenario -> distribution."""
        horizons = [
            (
                "30 Days",
                [
                    ("RM-100", ["Empirical", "Gamma-3"]),
                    ("RM-170", ["Empirical", "Gamma-3"]),
                    ("FFZ-350", ["Empirical", "Gamma-3"]),
                ],
            ),
            (
                "90 Days\n(Pareto-front policies only)",
                [
                    ("RM-100", ["Empirical", "Gamma-3"]),
                    ("RM-170", ["Empirical", "Gamma-3"]),
                    ("FFZ-350", ["Empirical", "Gamma-3"]),
                ],
            ),
        ]
        top = Inches(1.25)
        root_w, root_h = Inches(2.2), Inches(0.5)
        root_left = (SLIDE_W - root_w) / 2
        root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, root_left, top, root_w, root_h)
        _fill(root, DARK)
        root.shadow.inherit = False
        p = root.text_frame.paragraphs[0]
        p.text = "Simulation Runs"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        n_h = len(horizons)
        gap_h = Inches(0.4)
        hz_top = top + root_h + Inches(0.35)
        hz_w = int((SLIDE_W - Inches(1.2) - gap_h * (n_h - 1)) / n_h)
        hz_h = Inches(0.6)
        left0 = Inches(0.6)
        dist_h = Inches(0.4)
        dist_top = Inches(0.1)
        for hi, (hz_label, scenarios) in enumerate(horizons):
            hz_left = left0 + hi * (hz_w + gap_h)
            hz_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, hz_left, hz_top, hz_w, hz_h)
            _fill(hz_box, ACCENT)
            hz_box.shadow.inherit = False
            tf = hz_box.text_frame
            tf.word_wrap = True
            for li, line in enumerate(hz_label.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

            n_s = len(scenarios)
            gap_s = Inches(0.15)
            sc_top = hz_top + hz_h + Inches(0.3)
            sc_w = int((hz_w - gap_s * (n_s - 1)) / n_s)
            sc_h = Inches(0.5)
            for si, (sc_label, dists) in enumerate(scenarios):
                sc_left = hz_left + si * (sc_w + gap_s)
                sc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sc_left, sc_top, sc_w, sc_h)
                _fill(sc_box, MUTED)
                sc_box.shadow.inherit = False
                p = sc_box.text_frame.paragraphs[0]
                p.text = sc_label
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

                dist_top = sc_top + sc_h + Inches(0.1)
                dist_h = Inches(0.4)
                _, dtf = _textbox(slide, sc_left, dist_top, sc_w, dist_h)
                p = dtf.paragraphs[0]
                p.text = " / ".join(dists)
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER
        return dist_top + dist_h

    def _figure_slide(self, spec: dict) -> None:
        slide = self._new_slide()
        self._title_bar(slide, spec["title"])
        fig_dir = Path(spec["figures_dir"]) if spec.get("figures_dir") else self.figures_dir
        figures = spec.get("figures") or [spec.get("figure")]
        resolved = []
        for name in figures:
            p = fig_dir / name
            if not p.exists() and spec.get("fallback_figure"):
                p = fig_dir / spec["fallback_figure"]
            if p.exists():
                resolved.append(p)
            else:
                print(f"  [WARN] Figure not found: {p}")
        area_top = Inches(1.15)
        bottom_legend = spec.get("bottom_legend", False)
        legend_top = None
        legend_h = None
        if bottom_legend:
            fig_area_h = Inches(4.8)
            legend_top = area_top + fig_area_h + Inches(0.1)
            legend_h = SLIDE_H - legend_top - Inches(0.2)
        else:
            fig_area_h = SLIDE_H - area_top - Inches(0.2)
        area_h = fig_area_h

        legend_w = Inches(3.1) if spec.get("side_legend") else 0
        fig_area_w = SLIDE_W - Inches(0.6) - (legend_w + Inches(0.2) if legend_w else 0)
        if resolved:
            if spec.get("layout") == "vertical":
                h_each = int(fig_area_h / len(resolved))
                for i, p in enumerate(resolved):
                    self._picture_fit(slide, p, Inches(0.3), area_top + h_each * i, fig_area_w, h_each - Inches(0.1))
            else:
                w_each = int(fig_area_w / len(resolved))
                for i, p in enumerate(resolved):
                    self._picture_fit(slide, p, Inches(0.3) + w_each * i, area_top, w_each - Inches(0.2), fig_area_h)

        if legend_w:
            legend_left = Inches(0.3) + fig_area_w + Inches(0.2)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, legend_left, area_top, legend_w, area_h)  # pyrefly: ignore [bad-argument-type]
            _fill(box, RGBColor(0xF0, 0xF4, 0xFA))
            box.shadow.inherit = False
            _, tf = _textbox(
                slide,
                legend_left + Inches(0.2),
                area_top + Inches(0.15),
                legend_w - Inches(0.4),
                area_h - Inches(0.3),  # pyrefly: ignore [unsupported-operation]
            )
            p = tf.paragraphs[0]
            p.text = "Shared axis labels"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            for para in spec.get("side_legend_text", "").split("\n\n"):
                pp = tf.add_paragraph()
                pp.text = para
                pp.font.size = Pt(12)
                pp.font.color.rgb = DARK
                pp.space_before = Pt(10)
        elif bottom_legend:
            legend_left = Inches(0.3)
            legend_w = fig_area_w
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, legend_left, legend_top, legend_w, legend_h)  # pyrefly: ignore [bad-argument-type]
            _fill(box, RGBColor(0xF0, 0xF4, 0xFA))
            box.shadow.inherit = False
            _, tf = _textbox(
                slide,
                legend_left + Inches(0.2),
                legend_top + Inches(0.1),  # pyrefly: ignore [unsupported-operation]
                legend_w - Inches(0.4),
                legend_h - Inches(0.2),  # pyrefly: ignore [unsupported-operation]
            )
            p = tf.paragraphs[0]
            p.text = "Shared axis labels"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT

            legend_text = spec.get("side_legend_text") or spec.get("bottom_legend_text", "")
            for para in legend_text.split("\n\n"):
                pp = tf.add_paragraph()
                pp.text = para.replace("\n", " ")
                pp.font.size = Pt(10.5)
                pp.font.color.rgb = DARK
                pp.space_before = Pt(3)
        caption = spec.get("caption") or spec.get("note") or ""
        self._record_script(spec["title"], [caption])

    # ── Slides ──────────────────────────────────────────────────────────────────

    def cover(self) -> None:
        slide = self._new_slide()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)  # pyrefly: ignore [bad-argument-type]
        _fill(bg, DARK)
        title_len = len(self.content["title"])
        title_sz = 26 if title_len > 80 else 32 if title_len > 55 else 36
        band_top = Inches(4.55)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, band_top, SLIDE_W, Inches(0.06))  # pyrefly: ignore [bad-argument-type]
        _fill(band, ACCENT)
        _, tf = _textbox(slide, Inches(0.9), Inches(1.95), SLIDE_W - Inches(1.8), Inches(2.4))
        p = tf.paragraphs[0]
        p.text = self.content["title"]
        p.font.size = Pt(title_sz)
        p.font.bold = True
        p.font.color.rgb = WHITE
        _, tf2 = _textbox(slide, Inches(0.9), band_top + Inches(0.25), SLIDE_W - Inches(1.8), Inches(3.0))
        p3 = tf2.paragraphs[0]
        p3.text = self.author
        p3.font.size = Pt(19)
        p3.font.bold = True
        p3.font.color.rgb = WHITE
        if self.coauthors:
            p4 = tf2.add_paragraph()
            p4.text = "with " + ", ".join(self.coauthors)
            p4.font.size = Pt(14)
            p4.font.color.rgb = LIGHT_TXT
            p4.space_before = Pt(6)
        if self.groups:
            p5 = tf2.add_paragraph()
            p5.text = "   ·   ".join(self.groups)
            p5.font.size = Pt(11)
            p5.font.italic = True
            p5.font.color.rgb = RGBColor(0x8A, 0x9B, 0xB0)
            p5.space_before = Pt(10)
        # Institution logos at bottom; conference logo at top-right.
        # Bottom 3 logos (inescid, ist, cegist): equal horizontal spacing.
        _lw = [Inches(1.83), Inches(1.98), Inches(2.49)]
        _gap = int((SLIDE_W - sum(_lw)) / 4)
        _lx = [_gap, _gap + _lw[0] + _gap, _gap + _lw[0] + _gap + _lw[1] + _gap]
        logo_specs = [
            ("logo-inescid.png",          _lx[0], Inches(5.95), _lw[0], Inches(1.25)),
            ("logo-ist.png",              _lx[1], Inches(5.82), _lw[1], Inches(1.69)),
            ("logo-cegist.png",           _lx[2], Inches(5.97), _lw[2], Inches(1.38)),
            ("logo-optimization2026.png", Inches(10.39), Inches(0.15), Inches(2.77), Inches(1.35)),
        ]
        # Contain-fit each logo inside its [lx, ly, lw, lh] box at natural
        # aspect ratio so none of them renders stretched.
        for fname, lx, ly, lw, lh in logo_specs:
            logo_path = ASSETS_DIR / fname
            if not logo_path.exists():
                continue
            fit_x, fit_y, fit_w, fit_h = lx, ly, lw, lh
            try:
                from PIL import Image as _PILImage

                with _PILImage.open(logo_path) as im:
                    nat_w, nat_h = im.size
                if nat_w > 0 and nat_h > 0:
                    scale = min(int(lw) / nat_w, int(lh) / nat_h)
                    fit_w = int(nat_w * scale)
                    fit_h = int(nat_h * scale)
                    fit_x = int(lx) + (int(lw) - fit_w) // 2
                    fit_y = int(ly) + (int(lh) - fit_h) // 2
            except Exception:
                pass  # keep the box as-is if the image cannot be measured
            slide.shapes.add_picture(str(logo_path), fit_x, fit_y, fit_w, fit_h) # pyrefly: ignore [bad-argument-type]
        self._record_script(
            self.content["title"],
            [f"Presented by {self.author}." + (f" With {', '.join(self.coauthors)}." if self.coauthors else "")],
        )

    def agenda(self) -> None:
        slide = self._new_slide()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)  # pyrefly: ignore [bad-argument-type]
        _fill(bg, DARK)
        _, tf = _textbox(slide, Inches(0.9), Inches(0.5), SLIDE_W - Inches(1.8), Inches(0.9))
        p = tf.paragraphs[0]
        p.text = "Agenda"
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = WHITE

        items = self.content["agenda"]
        half = (len(items) + 1) // 2
        cols = [items[:half], items[half:]]
        top0 = Inches(1.65)
        avail_h = SLIDE_H - top0 - Inches(0.5)
        gap_x = Inches(0.4)
        col_w = int((SLIDE_W - Inches(1.2) - gap_x) / 2)
        card_colors = [ACCENT, RGBColor(0x3E, 0x8E, 0x41), RGBColor(0xB0, 0x6A, 0x2E), RGBColor(0x8E, 0x3E, 0x7A)]
        for col, chunk in enumerate(cols):
            left = Inches(0.6) + col * (col_w + gap_x)
            card_gap = Inches(0.22)
            card_h = int((avail_h - card_gap * (len(chunk) - 1)) / max(len(chunk), 1))
            for i, item in enumerate(chunk):
                idx = col * half + i
                top = top0 + i * (card_h + card_gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, card_h) # pyrefly: ignore [bad-argument-type]
                _fill(card, RGBColor(0x27, 0x37, 0x49))
                card.line.color.rgb = card_colors[idx % len(card_colors)] # pyrefly: ignore [bad-assignment]
                card.line.width = Pt(1.25)
                card.shadow.inherit = False
                badge_d = min(Inches(0.62), card_h - Inches(0.15))
                badge_left = left + Inches(0.18)
                badge_top = top + int((card_h - badge_d) / 2)
                badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_left, badge_top, badge_d, badge_d) # pyrefly: ignore [bad-argument-type]
                _fill(badge, card_colors[idx % len(card_colors)])
                badge.shadow.inherit = False
                bp = badge.text_frame.paragraphs[0]
                bp.text = str(idx + 1)
                bp.font.size = Pt(22)
                bp.font.bold = True
                bp.font.color.rgb = WHITE
                bp.alignment = PP_ALIGN.CENTER
                _, itf = _textbox(
                    slide, badge_left + badge_d + Inches(0.2), top, col_w - badge_d - Inches(0.55), card_h
                )
                itf.vertical_anchor = MSO_ANCHOR.MIDDLE
                ip = itf.paragraphs[0]
                ip.text = item
                ip.font.size = Pt(19)
                ip.font.bold = True
                ip.font.color.rgb = WHITE
        self._record_script("Agenda", ["Today's agenda: " + "; ".join(items) + "."])

    def _figures_side_by_side(self, slide, paths: list[Path], left, top, total_w, total_h) -> None:
        w_each = int(total_w / len(paths))
        for i, p in enumerate(paths):
            self._picture_fit(slide, p, left + w_each * i, top, w_each - Inches(0.1), total_h)

    def _figures_equal_height(self, slide, paths: list[Path], left, top, total_w, total_h) -> None:
        """Place images side by side sharing one common height (no per-image aspect-fit height drift)."""
        from PIL import Image

        aspects = []
        for p in paths:
            with Image.open(p) as im:
                w_px, h_px = im.size
            aspects.append(w_px / h_px)
        common_h = total_h
        total_needed_w = sum(a * common_h for a in aspects)
        if total_needed_w > total_w:
            common_h = int(common_h * (total_w / total_needed_w))
            total_needed_w = total_w
        group_left = left + int((total_w - total_needed_w) / 2)
        x = group_left
        for p, a in zip(paths, aspects, strict=True):
            w = int(a * common_h)
            slide.shapes.add_picture(str(p), x, top + int((total_h - common_h) / 2), w, common_h)
            x += w

    def _figures_grid_2x2(self, slide, paths: list[Path], left, top, total_w, total_h) -> None:
        gap = Inches(0.15)
        cell_w = int((total_w - gap) / 2)
        cell_h = int((total_h - gap) / 2)
        positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
        for p, (cx, cy) in zip(paths, positions, strict=False):
            px = left + cx * (cell_w + gap)
            py = top + cy * (cell_h + gap)
            self._picture_fit(slide, p, px, py, cell_w, cell_h)

    def content_slide(self, key: str) -> None:
        spec = self.content["slides"][key]
        slide = self._new_slide()
        self._title_bar(slide, spec["title"])
        show_bullets = spec.get("show_bullets", True)
        fig_names = spec.get("figures") or ([spec["figure"]] if spec.get("figure") else [])
        candidate_paths = [self.figures_dir / n for n in fig_names]
        fig_paths = []
        for p in candidate_paths:
            if p.exists():
                fig_paths.append(p)
            else:
                print(f"  [WARN] Figure not found: {p}")
        has_fig = bool(fig_paths)

        if spec.get("equation"):
            col_w = int(SLIDE_W / 2) - Inches(0.4) if has_fig else SLIDE_W - Inches(1.2)
            eq_left = Inches(0.6)
            content_bottom = self._equation_focus(
                slide, spec["equation"], left=eq_left, width=col_w,
                size_pt=spec.get("eq_size_pt", 15 if has_fig else 20), line_h=spec.get("eq_line_h", 0.62),
            )
            if spec.get("variables"):
                content_bottom = self._variable_glossary(
                    slide, spec["variables"], left=eq_left, top=content_bottom, width=col_w,
                    max_height=SLIDE_H - Inches(0.25) - content_bottom,
                    ncols=spec.get("variables_cols", 2), fontsize=spec.get("variables_fontsize", 12.5),
                )
            if show_bullets:
                self._bullets(
                    slide, spec["bullets"], left=eq_left, top=content_bottom, width=col_w, size=13 if has_fig else 14
                )
            if has_fig:
                self._figures_side_by_side(
                    slide, fig_paths, int(SLIDE_W / 2) + Inches(0.2), Inches(1.2),
                    int(SLIDE_W / 2) - Inches(0.8), SLIDE_H - Inches(1.5),
                )
        elif spec.get("diagram") == "pipeline":
            diagram_bottom = self._pipeline_diagram(slide)
            bullets_top = diagram_bottom + Inches(0.1)
            if show_bullets:
                self._bullets(slide, spec["bullets"], top=bullets_top, size=14)
        elif spec.get("diagram") == "policy_grid":
            self._policy_grid_diagram(slide)
            bullets_top = Inches(6.6)
            if show_bullets:
                self._bullets(slide, spec["bullets"], top=bullets_top, size=13)
        elif spec.get("diagram") == "metaheuristic_families":
            grid_bottom = SLIDE_H - Inches(0.3)
            if len(fig_paths) >= 3:
                self._figures_grid_2x2(slide, fig_paths, Inches(0.4), Inches(1.15), SLIDE_W - Inches(0.8), grid_bottom - Inches(1.15))
            else:
                self._figures_side_by_side(slide, fig_paths, Inches(0.4), Inches(1.15), SLIDE_W - Inches(0.8), grid_bottom - Inches(1.15))
            if show_bullets:
                self._bullets(slide, spec["bullets"], top=grid_bottom, size=12)
        elif spec.get("diagram") == "algo_taxonomy":
            self._algo_taxonomy_diagram(slide)
            bullets_top = Inches(4.35)
            bullets_w = int(SLIDE_W / 2) - Inches(0.7) if has_fig else None
            if show_bullets:
                self._bullets(slide, spec["bullets"], top=bullets_top, width=bullets_w, size=12, gap=5)
            if has_fig:
                self._figures_side_by_side(
                    slide, fig_paths, int(SLIDE_W / 2) + Inches(0.1), Inches(1.2),
                    int(SLIDE_W / 2) - Inches(0.7), SLIDE_H - Inches(1.5),
                )
        elif spec.get("diagram") == "doe_tree":
            diagram_bottom = self._doe_tree_diagram(slide)
            content_top = diagram_bottom + Inches(0.15)
            if fig_paths:
                self._figures_side_by_side(
                    slide, fig_paths, Inches(0.6), content_top,
                    SLIDE_W - Inches(1.2), SLIDE_H - content_top - Inches(0.2),
                )
            elif show_bullets:
                self._bullets(slide, spec["bullets"], top=content_top, size=13)
        elif has_fig:
            if show_bullets:
                fig_left, fig_w = int(SLIDE_W / 2), int(SLIDE_W / 2) - Inches(0.4)
            else:
                fig_left, fig_w = Inches(0.5), SLIDE_W - Inches(1.0)
            self._figures_equal_height(slide, fig_paths, fig_left, Inches(1.2), fig_w, SLIDE_H - Inches(1.4))
            if show_bullets:
                self._bullets(
                    slide, spec["bullets"], left=Inches(0.5), top=Inches(1.4),
                    width=int(SLIDE_W / 2) - Inches(0.7), size=13,
                )
        elif show_bullets:
            self._bullets(slide, spec["bullets"])
        script_parts = [spec.get("caption", "")] + list(spec["bullets"]) + list(spec.get("speaker_notes", []))
        self._record_script(spec["title"], script_parts)

    # ── Native-shape slides (slides 4 & 5) ─────────────────────────────────────

    def objective(self) -> None:
        """Slide 4: objective diagram + algorithm taxonomy (all native PPTX shapes, no image)."""
        spec = self.content["slides"]["objective"]
        slide = self._new_slide()
        I = 914400  # EMU per inch

        # ── Title bar ──
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))  # pyrefly: ignore [bad-argument-type]
        _fill(bar, DARK)
        bar.text_frame.word_wrap = True
        bar.text_frame.margin_left = Inches(0.5)
        p = bar.text_frame.paragraphs[0]
        p.text = spec["title"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # ── Subtitle band below title ──
        sub_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(4.126 * I), int(1.242 * I), int(5.634 * I), int(0.438 * I))  # pyrefly: ignore [bad-argument-type]
        _fill(sub_box, DARK)
        sub_box.line.fill.background()
        sub_box.shadow.inherit = False
        p = sub_box.text_frame.paragraphs[0]
        p.text = "One framework to compare them all"
        p.font.size = Pt(15)
        p.font.italic = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # ── Top flow diagram ──
        GREEN = RGBColor(0x3E, 0x8E, 0x41)
        ORANGE = RGBColor(0xB0, 0x6A, 0x2E)

        def _rbox(x_in, y_in, w_in, h_in, fill_color, text, font_sz=14):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x_in * I), int(y_in * I), int(w_in * I), int(h_in * I))  # pyrefly: ignore [bad-argument-type]
            box.fill.solid()
            box.fill.fore_color.rgb = fill_color
            box.line.fill.background()
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_sz)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            return box

        _rbox(2.059, 1.791, 1.426, 0.457, ACCENT, "Exact Methods", 13)
        _rbox(2.059, 2.414, 1.426, 0.600, GREEN, "Meta-Heuristics", 13)
        _rbox(2.059, 3.159, 1.426, 0.600, ORANGE, "Hyper-Heuristics", 13)

        # One Shared Simulator (dark box, two lines of text)
        sim_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(5.049 * I), int(1.704 * I), int(2.426 * I), int(2.046 * I))  # pyrefly: ignore [bad-argument-type]
        _fill(sim_box, DARK)
        sim_box.line.fill.background()
        sim_box.shadow.inherit = False
        sim_tf = sim_box.text_frame
        sim_tf.word_wrap = True
        sim_p1 = sim_tf.paragraphs[0]
        sim_p1.text = "One Shared Simulator"
        sim_p1.font.size = Pt(15)
        sim_p1.font.bold = True
        sim_p1.font.color.rgb = WHITE
        sim_p1.alignment = PP_ALIGN.CENTER
        sim_p2 = sim_tf.add_paragraph()
        sim_p2.text = "Multi-day simulator:\nregion × graph size × demand"
        sim_p2.font.size = Pt(11)
        sim_p2.font.color.rgb = LIGHT_TXT
        sim_p2.alignment = PP_ALIGN.CENTER
        sim_p2.space_before = Pt(6)

        # Fair Benchmarks (KPIs)
        _rbox(8.612, 2.158, 1.564, 1.139, ACCENT, "Fair Benchmarks\n(KPIs)", 13)

        # Connectors: Exact→Simulator (down-right, no flip)
        _add_connector(slide, 3.485 * I, 2.020 * I, 1.564 * I, 0.707 * I, "2E74B5")
        # Meta→Simulator (nearly horizontal, no flip)
        _add_connector(slide, 3.485 * I, 2.714 * I, 1.564 * I, 0.013 * I, "3E8E41")
        # Hyper→Simulator (up-right from bottom-left, flipV)
        _add_connector(slide, 3.485 * I, 2.727 * I, 1.564 * I, 0.732 * I, "B06A2E", flip_v=True)
        # Simulator→Benchmarks (horizontal, no flip)
        _add_connector(slide, 7.475 * I, 2.727 * I, 1.136 * I, 0.001 * I, "2E74B5")

        # ── Bottom taxonomy grid ──
        LIGHT_FILL = RGBColor(0xF0, 0xF4, 0xFA)

        def _tax_header(x_in, y_in, fill_color, text):
            hd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x_in * I), int(y_in * I), int(3.811 * I), int(0.600 * I))  # pyrefly: ignore [bad-argument-type]
            _fill(hd, fill_color)
            hd.shadow.inherit = False
            p = hd.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

        def _tax_item(x_in, y_in, h_in, text, line_color, font_sz=13):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x_in * I), int(y_in * I), int(3.811 * I), int(h_in * I))  # pyrefly: ignore [bad-argument-type]
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_FILL # pyrefly: ignore [bad-assignment]
            box.line.color.rgb = line_color
            box.line.width = Pt(1.25)
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_top = Pt(1)
            tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_sz)
            p.font.bold = True
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER

        # Exact Methods column (full algorithm names with acronyms)
        _tax_header(0.640, 4.258, ACCENT, "Exact Methods")
        _tax_item(0.640, 5.008, 1.015, "Branch-and-Price-and-Cut (BPC)", ACCENT)
        _tax_item(0.640, 6.143, 1.015, "Smart Waste Collection —\nTwo-Commodity Flow (SWC-TCF)", ACCENT)

        # Meta-Heuristics column — taller items so the long full names
        # (PG-CLNS, PSOMA) stay inside their boxes
        _tax_header(4.801, 4.258, GREEN, "Meta-Heuristics")
        _meta_algos = [
            "Hybrid Genetic Search (HGS)",
            "Adaptive Large Neighborhood Search (ALNS)",
            "Simulated Annealing Neighborhood Search (SANS)",
            "Pheromone-Guided Cooperative Large Neighborhood Search (PG-CLNS)",
            "Particle Swarm Optimization Memetic Algorithm (PSOMA)",
        ]
        for ii, algo in enumerate(_meta_algos):
            _tax_item(4.801, 5.008 + ii * 0.490, 0.450, algo, GREEN, font_sz=10)

        # Hyper-Heuristics column
        _tax_header(8.962, 4.258, ORANGE, "Hyper-Heuristics")
        _tax_item(8.962, 5.008, 2.150, "Ant Colony Optimization\nHyper-Heuristic (ACO-HH)", ORANGE)

        speaker_notes = spec.get("speaker_notes", [])
        self._record_script(spec["title"], speaker_notes)

    def simulator(self) -> None:
        """Slide 5: pipeline chevrons + bin/truck diagram (native shapes + images)."""
        spec = self.content["slides"]["simulator"]
        slide = self._new_slide()
        self._title_bar(slide, spec["title"])

        # Pipeline chevrons (re-use _pipeline_diagram, which no longer adds the simulator box)
        self._pipeline_diagram(slide)

        # ── Bin / truck diagram ──
        I = 914400
        BIN_IMG = Path(__file__).resolve().parent / "images" / "waste_bin_icon.png"
        TRUCK_IMG = Path(__file__).resolve().parent / "images" / "waste_truck_icon.png"
        DARK_NAVY = RGBColor(0x1F, 0x2D, 0x3D)
        RED = RGBColor(0xFF, 0x00, 0x00)

        def _bin_and_label(x_in, y_in, pct_text):
            """Place bin icon (0.495"×0.703") and percentage label overlaid near the bin bottom."""
            bw, bh = int(0.495 * I), int(0.703 * I)
            if BIN_IMG.exists():
                slide.shapes.add_picture(str(BIN_IMG), int(x_in * I), int(y_in * I), bw, bh)  # pyrefly: ignore [bad-argument-type]
            # Label sits in the lower portion of the bin icon (tmp/: label.y ≈ bin.y + 0.360")
            lx = int(x_in * I) + int(0.023 * I)
            ly = int((y_in + 0.360) * I)
            lw, lh = int(0.449 * I), int(0.286 * I)
            _, lbl_tf = _textbox(slide, lx, ly, lw, lh)
            lbl_tf.word_wrap = False
            lp = lbl_tf.paragraphs[0]
            lp.text = pct_text
            lp.font.size = Pt(11)
            lp.font.bold = True
            lp.font.color.rgb = WHITE
            lp.alignment = PP_ALIGN.CENTER

        # Left group (8 bins, scattered — before selection)
        _left_bins = [
            (1.063, 4.912, "87%"),
            (0.431, 5.399, "60%"),
            (0.708, 6.182, "42%"),
            (1.405, 5.600, "50%"),
            (1.538, 4.570, "30%"),
            (2.063, 5.113, "80%"),
            (2.412, 6.260, "20%"),
            (2.851, 5.335, "55%"),
        ]
        for bx, by, pct in _left_bins:
            _bin_and_label(bx, by, pct)

        # Red border around the "87%" selected bin in left group
        sel_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(0.990 * I), int(4.902 * I), int(0.599 * I), int(0.759 * I))  # pyrefly: ignore [bad-argument-type]
        sel_box.fill.background()
        sel_box.line.color.rgb = RED # pyrefly: ignore [bad-assignment]
        sel_box.line.width = Pt(2.5)
        sel_box.shadow.inherit = False

        # Right group (8 bins, same fill percentages — after selection)
        _right_bins = [
            (7.424, 4.698, "87%"),
            (6.791, 5.185, "60%"),
            (7.068, 5.968, "42%"),
            (7.766, 5.386, "50%"),
            (7.899, 4.356, "30%"),
            (8.424, 4.899, "80%"),
            (8.773, 6.046, "20%"),
            (9.212, 5.121, "55%"),
        ]
        for bx, by, pct in _right_bins:
            _bin_and_label(bx, by, pct)

        # Truck icon
        if TRUCK_IMG.exists():
            slide.shapes.add_picture(str(TRUCK_IMG), int(7.700 * I), int(6.673 * I), int(1.231 * I), int(0.943 * I)) # pyrefly: ignore [bad-argument-type]

        # Collection-tour arrows (orange): truck → 50 → 60 → 87 → 30 → 80 → truck.
        # Bounding boxes + flips copied verbatim from the reference deck
        # (tmp/wsmart_route_results.pptx slide 5) — the flips encode the tour
        # direction, so the arrowhead (tailEnd) lands on the tour target.
        _connector_coords = [
            (8.005, 6.032, 0.217, 0.756, True, True),   # truck → 50
            (7.245, 5.688, 0.528, 0.158, True, True),   # 50 → 60
            (7.039, 5.049, 0.385, 0.136, False, True),  # 60 → 87
            (7.706, 4.567, 0.226, 0.146, False, True),  # 87 → 30
            (8.299, 4.661, 0.373, 0.238, False, False),  # 30 → 80
            (8.462, 5.575, 0.140, 1.213, True, False),  # 80 → truck
        ]
        for cx0, cy0, ccx, ccy, fh, fv in _connector_coords:
            _add_connector(slide, cx0 * I, cy0 * I, ccx * I, ccy * I, "ED7D31", flip_v=fv, flip_h=fh, width_pt=1.5)

        # Horizontal underbrace beneath the constructor→improver chevrons — a
        # right brace rotated 90° (geometry from the reference deck; the
        # unrotated bbox spans vertically and rotates about its centre).
        brace = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_BRACE, int(7.977 * I), int(-0.194 * I), int(0.323 * I), int(8.752 * I) # pyrefly: ignore [bad-argument-type]
        )
        brace.rotation = 90
        brace.fill.background()
        brace.line.color.rgb = ACCENT # pyrefly: ignore [bad-assignment]
        brace.line.width = Pt(1.5)
        brace.shadow.inherit = False

        script_parts = list(spec.get("speaker_notes", [])) + list(spec["bullets"])
        self._record_script(spec["title"], script_parts)

    def _load_horizon(self, spec: dict) -> tuple | None:
        csv_path = Path(spec["csv"])
        if not csv_path.exists():
            return None
        df = filter_data(load_horizon_csv(csv_path), self._results_config)
        if df.empty:
            return None
        theme = {"name": "light", "fg": "#1F2D3D"}
        ctx = build_context(df, self._results_config, theme, spec["days"])
        return df, aggregate(df), ctx

    def _results_table_slide(self) -> None:
        """Full results table slide, for the horizon(s) chosen via --results-table."""
        if self.results_table == "none":
            return
        self._results_config = load_json("simulation_analysis_config.json")
        horizon_specs = sorted(self._results_config["horizons"], key=lambda h: h["days"])

        if self.results_table == "all":
            row_keys, col_keys, cells = [], [], {}
            used_days = []
            for spec in horizon_specs:
                loaded = self._load_horizon(spec)
                if loaded is None:
                    continue
                df_full, _, ctx = loaded
                df_cls = df_full[df_full["improver"].str.upper() == "CLS"]
                if df_cls.empty:
                    df_cls = df_full
                dfm_cls = aggregate(df_cls)
                ctx_cls = {**ctx, "improvers": list(df_cls["improver"].unique())}
                rk, ck, cs = build_full_results_matrix(dfm_cls, ctx_cls, horizon_label=f"{spec['days']}d")
                row_keys += [r for r in rk if r not in row_keys]
                col_keys += ck
                cells.update(cs)
                used_days.append(spec["days"])
            row_keys.sort()
            if not row_keys:
                print("  [WARN] No data available for the full results table — skipping slide")
                return
            title = f"Results — CLS Improver Table (All Horizons: {', '.join(f'{d}d' for d in used_days)})"
            multi = True
        else:
            days = int(self.results_table.rstrip("d"))
            spec = next((s for s in horizon_specs if s["days"] == days), None)
            loaded = self._load_horizon(spec) if spec else None
            if loaded is None:
                print(f"  [WARN] No data for the {self.results_table} horizon — skipping full results table slide")
                return
            df_full, _, ctx = loaded
            df_cls = df_full[df_full["improver"].str.upper() == "CLS"]
            if df_cls.empty:
                df_cls = df_full
            dfm_cls = aggregate(df_cls)
            ctx_cls = {**ctx, "improvers": list(df_cls["improver"].unique())}
            row_keys, col_keys, cells = build_full_results_matrix(dfm_cls, ctx_cls)
            title = f"Results — CLS Improver ({days}-Day Horizon)"
            multi = False

        level_phrases = {
            "horizon": "horizon",
            "strategy": "mandatory selection strategy",
            "constructor": "route constructor",
            "improver": "route improver",
        }
        level_names = (["horizon"] if multi else []) + ["strategy", "constructor", "improver"]
        col_desc = " × ".join(level_phrases[n] for n in level_names)
        row_labels = ["Region", "N", "Dist"]
        corner_note = "TOP:\nOverflows\n\nBOTTOM:\nKG / KM"

        slide = self._new_slide()
        self._title_bar(slide, title)
        area_top, area_bottom = Inches(1.15), SLIDE_H - Inches(0.1)
        area_w_in = Emu(SLIDE_W).inches
        area_h_in = Emu(area_bottom - area_top).inches

        global_best = compute_global_best(row_keys, col_keys, cells)
        img_path = render_hier_table_image(
            row_keys, col_keys, cells, row_labels, self._tmp / "full_results_table.png",
            target_size_in=(area_w_in, area_h_in), global_best=global_best, corner_note=corner_note,
        )
        slide.shapes.add_picture(str(img_path), 0, area_top, SLIDE_W, area_bottom - area_top)  # pyrefly: ignore [bad-argument-type]

        self._record_script(
            title,
            [f"CLS route improver results table, columns grouped by {col_desc}. "
             f"In each cell, top=overflows (lower is better), bottom={KGKM_LABEL} (higher is better); best in bold green."],
        )

    def acknowledgments(self) -> None:
        slide = self._new_slide()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)  # pyrefly: ignore [bad-argument-type]
        _fill(bg, DARK)
        _, tf = _textbox(slide, Inches(1.0), Inches(1.8), SLIDE_W - Inches(2.0), Inches(1.0))
        p = tf.paragraphs[0]
        p.text = "Acknowledgements"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        ack_text = (
            "This work was supported by FCT – Foundation for Science and Technology, I.P., "
            "under project 2022.04180.PTDC and individual PhD research grant - 2025.06860.BDANA"
        )
        _, tf2 = _textbox(slide, Inches(1.0), Inches(3.0), SLIDE_W - Inches(2.0), Inches(1.8))
        pp = tf2.paragraphs[0]
        pp.text = ack_text
        pp.font.size = Pt(17)
        pp.font.color.rgb = LIGHT_TXT
        pp.alignment = PP_ALIGN.CENTER
        # FCT logo (white version) centred in lower half
        fct_logo = Path(__file__).resolve().parent / "images" / "2022_FCT_Logo_B_horizontal_branco.png"
        if fct_logo.exists():
            from PIL import Image as _PIL_Image
            with _PIL_Image.open(fct_logo) as _im:
                _lw_px, _lh_px = _im.size
            _logo_w = Inches(4.5)
            _logo_h = int(_logo_w * _lh_px / _lw_px)
            _logo_x = int((SLIDE_W - _logo_w) / 2)
            _logo_y = Inches(5.0)
            slide.shapes.add_picture(str(fct_logo), _logo_x, _logo_y, _logo_w, _logo_h) # pyrefly: ignore [bad-argument-type]
        self._record_script("Acknowledgements", [ack_text])

    def qa(self) -> None:
        spec = self.content["slides"]["qa"]
        slide = self._new_slide()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)  # pyrefly: ignore [bad-argument-type]
        _fill(bg, DARK)
        fig_path = self.figures_dir / spec["figure"] if spec.get("figure") else None
        has_fig = bool(fig_path) and fig_path.exists()
        text_left = Inches(0.9)
        text_w = (int(SLIDE_W / 2) - Inches(1.1)) if has_fig else (SLIDE_W - Inches(1.8))
        _, tf = _textbox(slide, text_left, Inches(2.8), text_w, Inches(2.4))
        p = tf.paragraphs[0]
        p.text = spec["title"]
        p.font.size = Pt(30 if has_fig else 36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER if not has_fig else PP_ALIGN.LEFT
        for line in spec["bullets"]:
            pp = tf.add_paragraph()
            pp.text = line
            pp.font.size = Pt(16)
            pp.font.color.rgb = LIGHT_TXT
            pp.alignment = PP_ALIGN.CENTER if not has_fig else PP_ALIGN.LEFT
        if has_fig:
            self._picture_fit(
                slide,
                fig_path,
                int(SLIDE_W / 2) + Inches(0.2),
                Inches(0.6),
                int(SLIDE_W / 2) - Inches(0.8),
                SLIDE_H - Inches(1.2),
            )
        self._record_script(spec["title"], list(spec["bullets"]))

    def build(self) -> PresentationClass:
        self.cover()  # 1
        self.agenda()  # 2
        self.content_slide("vrpp")  # 3
        self.objective()  # 4
        self.simulator()  # 5
        self.content_slide("policy_overview")  # 6
        self.content_slide("strategies")  # 7
        self.content_slide("exact")  # 8
        self.content_slide("metaheuristics")  # 9
        self.content_slide("improvers")  # 10
        self.content_slide("design_of_experiments")  # 11
        self._figure_slide(self.content["figure_slides"]["pareto"])  # 12
        self._figure_slide(self.content["figure_slides"]["strategy_bubble"])  # 13
        self._figure_slide(self.content["figure_slides"]["scenario_heatmaps"])  # 14 (30d)
        self._figure_slide(self.content["figure_slides"]["heatmaps"])  # 15 (90d)
        self._figure_slide(self.content["figure_slides"]["improver_bubble"])  # 16
        self._results_table_slide()  # 17 (full results table, if requested)
        self.content_slide("conclusion")  # 18
        self.acknowledgments()  # 19
        self.qa()  # 20
        return self.prs


def generate_qa_route_image(out_path: Path) -> Path:
    """Generate a VRP with Profits routing illustration with multiple routes and unvisited nodes."""
    # Set seed for reproducibility
    np.random.seed(42)

    # Generate nodes
    n_nodes = 45
    coords = np.random.rand(n_nodes, 2)
    depot = np.array([0.5, 0.5])

    # Calculate angles and distances from depot
    diffs = coords - depot
    angles = np.arctan2(diffs[:, 1], diffs[:, 0])
    distances = np.linalg.norm(diffs, axis=1)

    route1_idx = []
    route2_idx = []
    route3_idx = []
    unvisited_idx = []

    for i in range(n_nodes):
        # Leave about 30% of the nodes unvisited to show the profit-based node selection
        if i % 3 == 0 and distances[i] > 0.15:
            unvisited_idx.append(i)
            continue

        angle = angles[i]
        if -np.pi <= angle < -np.pi / 3:
            route1_idx.append(i)
        elif -np.pi / 3 <= angle < np.pi / 3:
            route2_idx.append(i)
        else:
            route3_idx.append(i)

    # Sort nodes in each route by polar angle to make smooth tours
    def sort_route(indices):
        if not indices:
            return []
        sub_angles = angles[indices]
        sorted_indices = [indices[idx] for idx in np.argsort(sub_angles)]
        return sorted_indices

    r1 = sort_route(route1_idx)
    r2 = sort_route(route2_idx)
    r3 = sort_route(route3_idx)

    fig, ax = plt.subplots(figsize=(8, 8), dpi=200)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    # Draw faint gray edges between nearby nodes to represent the road network
    for i in range(n_nodes):
        for j in range(i + 1, n_nodes):
            dist = np.linalg.norm(coords[i] - coords[j])
            if dist < 0.22:
                ax.plot([coords[i, 0], coords[j, 0]], [coords[i, 1], coords[j, 1]], color="#E2E8F0", lw=0.6, zorder=1)

    # Plot unvisited nodes (selection part of VRPP)
    ax.scatter(
        coords[unvisited_idx, 0],
        coords[unvisited_idx, 1],
        color="#CBD5E1",
        edgecolor="#64748B",
        s=100,
        label="Unvisited Bins (No Profit)",
        zorder=3,
    )

    # Premium colors for the 3 routes (different vehicles)
    colors = ["#10B981", "#6366F1", "#F59E0B"]
    route_labels = ["Vehicle Route 1", "Vehicle Route 2", "Vehicle Route 3"]

    # Plot each route
    def plot_tour(route_nodes, color, label):
        if not route_nodes:
            return
        tour_coords = [depot] + [coords[idx] for idx in route_nodes] + [depot]
        tour_coords = np.array(tour_coords)

        ax.plot(tour_coords[:, 0], tour_coords[:, 1], color=color, lw=3, label=label, zorder=2)
        ax.scatter(tour_coords[1:-1, 0], tour_coords[1:-1, 1], color=color, edgecolor="black", s=120, zorder=4)

    plot_tour(r1, colors[0], route_labels[0])
    plot_tour(r2, colors[1], route_labels[1])
    plot_tour(r3, colors[2], route_labels[2])

    # Plot depot clearly marked
    ax.scatter(
        depot[0],
        depot[1],
        color="#EF4444",
        marker="s",
        s=220,
        edgecolor="black",
        linewidth=2,
        label="Depot / Warehouse",
        zorder=5,
    )

    ax.set_xlim(-0.05, 1.05)
    ax.set_ylim(-0.05, 1.05)
    ax.axis("off")

    ax.legend(
        loc="lower center",
        bbox_to_anchor=(0.5, -0.08),
        ncol=2,
        frameon=True,
        facecolor="white",
        edgecolor="#E2E8F0",
        fontsize=10,
    )

    plt.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


# ── Native reference-style diagrams (see links/reference_image_links.xml for
# the third-party images that inspired their layout/content) ──────────────────


def generate_bb_tree_image(out_path: Path) -> Path:
    """A small Branch-and-Bound search tree: BPC/SWC-TCF both branch, price and cut at each node."""
    fig, ax = plt.subplots(figsize=(8.1, 6), dpi=200)
    ax.set_xlim(0, 11.6)
    ax.set_ylim(0, 10)
    ax.axis("off")

    nodes = {
        "root": (5.6, 9, "LP Relaxation\n(+ cuts)", "#1F2D3D"),
        "l1": (2.3, 6.2, "Branch: x_ij = 0", "#2E74B5"),
        "r1": (8.6, 6.2, "Branch: x_ij = 1", "#2E74B5"),
        "l2": (2.3, 3.2, "Pruned\n(bound ≤ incumbent)", "#8A9BB0"),
        "r2a": (6.9, 3.2, "Branch again", "#2E74B5"),
        "r2b": (10.2, 3.2, "Integer feasible\nnew incumbent", "#3E8E41"),
        "leaf1": (5.3, 0.6, "Pruned\n(infeasible)", "#8A9BB0"),
        "leaf2": (8.7, 0.6, "Optimal route set", "#3E8E41"),
    }
    edges = [("root", "l1"), ("root", "r1"), ("l1", "l2"), ("r1", "r2a"), ("r1", "r2b"), ("r2a", "leaf1"), ("r2a", "leaf2")]
    for a, b in edges:
        xa, ya, *_ = nodes[a]
        xb, yb, *_ = nodes[b]
        ax.plot([xa, xb], [ya - 0.35, yb + 0.35], color="#8A9BB0", linewidth=1.4, zorder=1)

    for _key, (x, y, label, color) in nodes.items():
        w, h = 2.6, 0.9
        box = mpatches.FancyBboxPatch(
            (x - w / 2, y - h / 2), w, h, boxstyle="round,pad=0.06,rounding_size=0.12",
            facecolor=color, edgecolor="white", linewidth=1.2, zorder=2,
        )
        ax.add_patch(box)
        ax.text(x, y, label, ha="center", va="center", fontsize=9.5, color="white", fontweight="bold", zorder=3)

    ax.text(5.6, 9.85, "Branch-and-Bound (BPC / SWC-TCF)", ha="center", va="center", fontsize=13, fontweight="bold")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def _ls_axis(ax) -> None:
    ax.set_xlim(-0.15, 1.05)
    ax.set_ylim(-0.1, 1.05)
    ax.axis("off")


def _ls_points(ax, pts, labels, colors=None) -> None:
    colors = colors or ["#1F2D3D"] * len(pts)
    ax.scatter(pts[:, 0], pts[:, 1], s=150, c=colors, zorder=5, edgecolor="white", linewidth=1.5)
    for (x, y), lbl in zip(pts, labels, strict=True):
        ax.text(x, y + 0.1, lbl, ha="center", va="center", fontsize=11, fontweight="bold", color="#1F2D3D")


def generate_ls_operators_image(out_path: Path) -> Path:
    """Before/after illustration of three CLS neighbourhood moves: 2-opt, swap and relocate."""
    fig, axes = plt.subplots(2, 3, figsize=(13.5, 7.2), dpi=200)

    # -- 2-opt: remove a crossing by reversing the segment between the two edges --
    pts = np.array([[0.1, 0.15], [0.85, 0.75], [0.75, 0.1], [0.15, 0.85]])
    labels = ["A", "B", "C", "D"]
    ax = axes[0, 0]
    _ls_points(ax, pts, labels)
    ax.plot(*zip(pts[0], pts[1]), color="#C0392B", linewidth=3, zorder=3)
    ax.plot(*zip(pts[2], pts[3]), color="#C0392B", linewidth=3, zorder=3)
    ax.plot(*zip(pts[1], pts[2]), "--", color="#8A9BB0", linewidth=1.5, zorder=2)
    ax.plot(*zip(pts[3], pts[0]), "--", color="#8A9BB0", linewidth=1.5, zorder=2)
    _ls_axis(ax)
    ax.set_title("2-opt — before", fontsize=12, fontweight="bold", color="#C0392B")
    ax = axes[1, 0]
    _ls_points(ax, pts, labels)
    ax.plot(*zip(pts[0], pts[2]), color="#3E8E41", linewidth=3, zorder=3)
    ax.plot(*zip(pts[1], pts[3]), color="#3E8E41", linewidth=3, zorder=3)
    ax.plot(*zip(pts[1], pts[2]), "--", color="#8A9BB0", linewidth=1.5, zorder=2)
    ax.plot(*zip(pts[3], pts[0]), "--", color="#8A9BB0", linewidth=1.5, zorder=2)
    _ls_axis(ax)
    ax.set_title("2-opt — after: A–C, B–D", fontsize=12, fontweight="bold", color="#3E8E41")

    # -- Swap: exchange two nodes between routes --
    r1 = np.array([[0.05, 0.85], [0.4, 0.85], [0.75, 0.85]])
    r2 = np.array([[0.05, 0.15], [0.4, 0.15], [0.75, 0.15]])
    ax = axes[0, 1]
    _ls_points(ax, np.vstack([r1, r2]), ["A", "B", "C", "D", "E", "F"],
               colors=["#2E74B5"] * 3 + ["#B06A2E"] * 3)
    ax.plot(r1[:, 0], r1[:, 1], color="#2E74B5", linewidth=3, zorder=3)
    ax.plot(r2[:, 0], r2[:, 1], color="#B06A2E", linewidth=3, zorder=3)
    _ls_axis(ax)
    ax.set_title("Swap — before", fontsize=12, fontweight="bold", color="#C0392B")
    ax = axes[1, 1]
    r1s, r2s = r1.copy(), r2.copy()
    r1s[1], r2s[1] = r2[1], r1[1]
    _ls_points(ax, np.vstack([r1s, r2s]), ["A", "E", "C", "D", "B", "F"],
               colors=["#2E74B5"] * 3 + ["#B06A2E"] * 3)
    ax.plot(r1s[:, 0], r1s[:, 1], color="#2E74B5", linewidth=3, zorder=3)
    ax.plot(r2s[:, 0], r2s[:, 1], color="#B06A2E", linewidth=3, zorder=3)
    _ls_axis(ax)
    ax.set_title("Swap — after: B ↔ E", fontsize=12, fontweight="bold", color="#3E8E41")

    # -- Relocate: move a single node to a better position in the same route --
    before = np.array([[0.05, 0.5], [0.35, 0.85], [0.65, 0.15], [0.95, 0.5]])
    ax = axes[0, 2]
    _ls_points(ax, before, ["A", "B", "C", "D"])
    ax.plot(before[:, 0], before[:, 1], color="#C0392B", linewidth=3, zorder=3)
    _ls_axis(ax)
    ax.set_title("Relocate — before", fontsize=12, fontweight="bold", color="#C0392B")
    ax = axes[1, 2]
    after = np.array([[0.05, 0.5], [0.65, 0.15], [0.35, 0.85], [0.95, 0.5]])
    _ls_points(ax, after, ["A", "C", "B", "D"])
    ax.plot(after[:, 0], after[:, 1], color="#3E8E41", linewidth=3, zorder=3)
    _ls_axis(ax)
    ax.set_title("Relocate — after: B moved", fontsize=12, fontweight="bold", color="#3E8E41")

    fig.suptitle("Classical Local Search Neighbourhood Moves", fontsize=15, fontweight="bold")
    fig.tight_layout(rect=(0, 0, 1, 0.94))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_knapsack_image(out_path: Path) -> Path:
    """Illustrates mandatory selection as a knapsack decision: which bins fit today's budget."""
    fig, ax = plt.subplots(figsize=(7.4, 6.6), dpi=200)
    ax.set_xlim(0, 11.6)
    ax.set_ylim(0, 10)
    ax.axis("off")
    ax.set_title("Mandatory Selection ≈ a Knapsack Problem", fontsize=13, fontweight="bold")

    # Knapsack outline
    sack = mpatches.FancyBboxPatch(
        (0.6, 1.6), 4.0, 5.2, boxstyle="round,pad=0.05,rounding_size=0.3",
        facecolor="#F0F4FA", edgecolor="#1F2D3D", linewidth=2.5,
    )
    ax.add_patch(sack)
    ax.text(2.6, 7.55, "Today's Route\n(capacity Q)", ha="center", va="center", fontsize=11, fontweight="bold")

    items_in = [("Bin 1\n90% full", "#3E8E41"), ("Bin 2\n85% full", "#3E8E41"), ("Bin 3\n70% full", "#3E8E41")]
    for i, (label, color) in enumerate(items_in):
        y = 2.3 + i * 1.5
        box = mpatches.FancyBboxPatch((1.0, y), 3.2, 1.1, boxstyle="round,pad=0.03,rounding_size=0.12",
                                       facecolor=color, edgecolor="white", linewidth=1.2)
        ax.add_patch(box)
        ax.text(2.6, y + 0.55, label, ha="center", va="center", fontsize=10, color="white", fontweight="bold")

    items_out = ["Bin 4\n30% full", "Bin 5\n20% full", "Bin 6\n15% full"]
    for i, label in enumerate(items_out):
        x = 6.4 + (i % 2) * 2.5
        y = 7.4 - (i // 2) * 1.6
        box = mpatches.FancyBboxPatch((x, y), 2.1, 1.1, boxstyle="round,pad=0.03,rounding_size=0.12",
                                       facecolor="#CBD5E1", edgecolor="#5A6A7A", linewidth=1.2)
        ax.add_patch(box)
        ax.text(x + 1.05, y + 0.55, label, ha="center", va="center", fontsize=9.5, color="#333333", fontweight="bold")

    ax.text(
        5.8, 0.9,
        "maximise Σ value (waste kg)  subject to  Σ weight (fill/urgency) ≤ today's budget",
        ha="center", va="center", fontsize=10.5, fontstyle="italic", color="#5A6A7A",
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_lookahead_flow_image(out_path: Path) -> Path:
    """The Look-Ahead (LA) mandatory-selection strategy's multi-phase decision flow
    (ports `LookaheadSelection.select_bins`, logic/src/policies/mandatory_selection/
    selection_lookahead.py): flag today's overflowing bins, simulate their collection,
    find the next day any of them would overflow again, and bundle in any other bin
    that would overflow before that day so it isn't visited on its own extra trip."""
    DARK, ACCENT, GREEN, GREY = "#1F2D3D", "#2E74B5", "#3E8E41", "#5A6A7A"

    fig, ax = plt.subplots(figsize=(6.9, 10.2), dpi=200)
    ax.set_xlim(-0.3, 11.3)
    ax.set_ylim(-0.4, 17.0)
    ax.axis("off")
    ax.text(5, 16.5, "Look-Ahead (LA) — multi-phase decision flow", ha="center", va="center",
             fontsize=12.5, fontweight="bold", color=DARK)

    def _box(cx, cy, w, h, text, color, fontsize=9.5):
        b = mpatches.FancyBboxPatch(
            (cx - w / 2, cy - h / 2), w, h, boxstyle="round,pad=0.05,rounding_size=0.12",
            facecolor=color, edgecolor="white", linewidth=1.2, zorder=2,
        )
        ax.add_patch(b)
        ax.text(cx, cy, text, ha="center", va="center", fontsize=fontsize, color="white",
                 fontweight="bold", zorder=3)

    def _diamond(cx, cy, w, h, text, fontsize=9.5):
        pts = [(cx, cy + h / 2), (cx + w / 2, cy), (cx, cy - h / 2), (cx - w / 2, cy)]
        ax.add_patch(mpatches.Polygon(pts, closed=True, facecolor=ACCENT, edgecolor="white",
                                        linewidth=1.2, zorder=2))
        ax.text(cx, cy, text, ha="center", va="center", fontsize=fontsize, color="white",
                 fontweight="bold", zorder=3)

    def _arrow(x1, y1, x2, y2, color=GREY, lw=1.6):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1), arrowprops=dict(arrowstyle="-|>", color=color, linewidth=lw))

    def _elbow(x1, y1, y2, x2, color=GREY, lw=1.6, label=None):
        """Route right, then vertically, then arrow left/right into (x2, y2) — the flowchart's loop-back edges."""
        x_mid = x1 + 0.6
        ax.plot([x1, x_mid], [y1, y1], color=color, linewidth=lw, zorder=1)
        ax.plot([x_mid, x_mid], [y1, y2], color=color, linewidth=lw, zorder=1)
        _arrow(x_mid, y2, x2, y2, color=color, lw=lw)
        if label:
            ax.text(x_mid + 0.15, (y1 + y2) / 2, label, fontsize=8.5, color=color, ha="left", va="center")

    # Step 1: flag mandatory bins
    _box(5, 15.4, 8.6, 1.3, r"Day $t$: retrieve $w_{i,t}$" "\n" r"compute $w_{i,t}+\bar{\delta}_i,\ \forall i$",
         DARK, fontsize=10)
    _arrow(5, 14.75, 5, 14.0)
    _diamond(5, 13.0, 7.0, 2.0, r"$\exists\, i: w_{i,t}+\bar{\delta}_i \geq \tau$ ?", fontsize=9.5)
    _box(8.9, 13.0, 1.7, 1.0, r"$t=t{+}1$", GREY, fontsize=9.5)
    _arrow(8.5, 13.0, 8.75, 13.0, color=GREY)
    ax.text(8.15, 13.55, "No", fontsize=8.5, color=GREY, fontweight="bold")
    _elbow(9.75, 13.0, 15.4, 8.75, color=GREY)  # t=t+1 loops back up to the top box
    ax.text(5, 12.0 - 0.35, "Yes", fontsize=8.5, color=DARK, fontweight="bold", ha="center")
    _arrow(5, 12.0, 5, 11.2)

    # Step 2 (empty M_0 -> stop) is the diamond's "No" branch above: nothing flagged, day repeats.
    # Step 3: simulate collection of M_0
    _box(5, 10.4, 8.8, 1.6,
         r"$\mathcal{M}_0=\{\,i : w_{i,t}+\bar{\delta}_i \geq \tau\,\}$" "\n"
         r"simulate collection: $\hat{w}_{i,t}=0,\ \forall i\in\mathcal{M}_0$",
         DARK, fontsize=9.5)
    _arrow(5, 9.6, 5, 9.3)

    # Step 4: next collection day nc — simulate M_0 forward day by day
    _box(5, 8.85, 4.6, 0.9, r"$t'=t+1$", ACCENT, fontsize=10)
    _arrow(5, 8.4, 5, 8.05)
    _box(5, 7.55, 7.4, 1.0, r"compute $\hat{w}_{i,t'},\ \forall i\in\mathcal{M}_0$", ACCENT, fontsize=9.5)
    _arrow(5, 7.05, 5, 6.5)
    _diamond(5, 5.55, 7.2, 1.9, r"$\exists\, i\in\mathcal{M}_0: \hat{w}_{i,t'} \geq \tau$ ?", fontsize=9.2)
    _box(8.9, 5.55, 1.7, 1.0, r"$t'=t'{+}1$", GREY, fontsize=9.5)
    _arrow(8.6, 5.55, 8.75, 5.55, color=GREY)
    ax.text(8.15, 6.1, "No", fontsize=8.5, color=GREY, fontweight="bold")
    _elbow(9.75, 5.55, 7.55, 8.7, color=GREY, label="recompute")  # loops back to the compute box
    ax.text(5, 4.6 - 0.35, "Yes", fontsize=8.5, color=DARK, fontweight="bold", ha="center")
    _arrow(5, 4.6, 5, 3.85)

    # nc found
    _box(5, 3.35, 6.0, 1.0, r"$nc=t'$   (next collection day)", DARK, fontsize=10)
    _arrow(5, 2.85, 5, 2.5)

    # Step 5/6: bundle bins that would overflow before nc
    _box(5, 1.85, 9.0, 1.3,
         r"$\forall i \notin \mathcal{M}_0,\ t''=t{+}1,\dots,nc{-}1:$" "\n"
         r"check $w_{i,t}+(t''-t)\bar{\delta}_i \geq \tau$",
         ACCENT, fontsize=9.5)
    _arrow(5, 1.2, 5, 0.9)
    _box(5, 0.4, 9.0, 1.0,
         r"$\mathcal{M}=\mathcal{M}_0\ \cup\ \{\,i: \text{overflows before } nc\,\}$",
         GREEN, fontsize=10)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_framework_objective_image(out_path: Path) -> Path:
    """One simulator, many algorithms in, one comparable benchmark out."""
    fig, ax = plt.subplots(figsize=(6.2, 6.6), dpi=200)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")

    algos = ["Exact\nMethods", "Meta-\nHeuristics", "Hyper-\nHeuristics"]
    colors = [ACCENT_HEX, "#3E8E41", "#B06A2E"]
    for i, (label, color) in enumerate(zip(algos, colors, strict=True)):
        y = 8.4 - i * 1.7
        box = mpatches.FancyBboxPatch((0.4, y), 2.6, 1.15, boxstyle="round,pad=0.04,rounding_size=0.14",
                                       facecolor=color, edgecolor="white", linewidth=1.4)
        ax.add_patch(box)
        ax.text(1.7, y + 0.575, label, ha="center", va="center", fontsize=11, color="white", fontweight="bold")
        ax.annotate("", xy=(4.4, 5.2), xytext=(3.0, y + 0.575),
                    arrowprops=dict(arrowstyle="-|>", color="#8A9BB0", linewidth=1.6))

    sim = mpatches.FancyBboxPatch((4.4, 4.2), 3.0, 2.0, boxstyle="round,pad=0.05,rounding_size=0.2",
                                   facecolor="#1F2D3D", edgecolor="white", linewidth=1.6)
    ax.add_patch(sim)
    ax.text(5.9, 5.2, "One Shared\nSimulator", ha="center", va="center", fontsize=12, color="white",
            fontweight="bold")

    ax.annotate("", xy=(7.9, 5.2), xytext=(7.4, 5.2),
                arrowprops=dict(arrowstyle="-|>", color="#8A9BB0", linewidth=2.0))
    out = mpatches.FancyBboxPatch((7.9, 4.2), 1.7, 2.0, boxstyle="round,pad=0.05,rounding_size=0.2",
                                   facecolor="#5A6A7A", edgecolor="white", linewidth=1.4)
    ax.add_patch(out)
    ax.text(8.75, 5.2, "Fair\nBenchmark", ha="center", va="center", fontsize=10.5, color="white",
            fontweight="bold")

    ax.text(
        5.9, 1.6,
        "Same scenarios, same KPIs (overflows, KG / KM)\nfor every exact / meta- / hyper-heuristic method",
        ha="center", va="center", fontsize=10.5, fontstyle="italic", color="#5A6A7A",
    )
    ax.set_title("Objective: One Framework to Compare Them All", fontsize=13, fontweight="bold")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_metaheuristic_overview_image(out_path: Path) -> Path:
    """A fitness-landscape view of meta-heuristic explore/exploit search."""
    fig, ax = plt.subplots(figsize=(8.6, 4.0), dpi=200)
    x = np.linspace(0, 10, 400)
    y = (
        3 + 1.4 * np.sin(x * 1.3) + 0.6 * np.sin(x * 3.7 + 1) - 0.12 * (x - 5) ** 2 * 0.15
    )
    ax.plot(x, y, color="#1F2D3D", linewidth=2.5)
    ax.fill_between(x, y, y.min() - 0.5, color="#F0F4FA")
    # incumbent walking down toward a local/global optimum
    idx_start, idx_mid, idx_end = 40, 180, 305
    ax.scatter([x[idx_start]], [y[idx_start]], s=140, color="#B06A2E", zorder=5, edgecolor="white", linewidth=1.5)
    ax.annotate("start", (x[idx_start], y[idx_start]), textcoords="offset points", xytext=(0, 14),
                ha="center", fontsize=10, fontweight="bold", color="#B06A2E")
    ax.annotate(
        "", xy=(x[idx_mid], y[idx_mid] + 0.15), xytext=(x[idx_start], y[idx_start] + 0.15),
        arrowprops=dict(arrowstyle="-|>", color="#C0392B", linewidth=1.8, connectionstyle="arc3,rad=-0.3"),
    )
    ax.text((x[idx_start] + x[idx_mid]) / 2, y[idx_mid] + 1.1, "explore\n(escape local optima)",
            ha="center", fontsize=9.5, color="#C0392B", fontweight="bold")
    ax.scatter([x[idx_end]], [y[idx_end]], s=170, color="#3E8E41", zorder=5, edgecolor="white", linewidth=1.5)
    ax.annotate(
        "", xy=(x[idx_end], y[idx_end] + 0.1), xytext=(x[idx_mid], y[idx_mid] + 0.1),
        arrowprops=dict(arrowstyle="-|>", color="#3E8E41", linewidth=1.8, connectionstyle="arc3,rad=0.2"),
    )
    ax.text((x[idx_mid] + x[idx_end]) / 2, y[idx_end] - 1.0, "exploit\n(intensify the best region)",
            ha="center", fontsize=9.5, color="#3E8E41", fontweight="bold")
    ax.set_title("Meta-Heuristics: Explore the Landscape, Exploit the Best Region", fontsize=12, fontweight="bold")
    ax.axis("off")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_hyperheuristic_overview_image(out_path: Path) -> Path:
    """A controller selecting among low-level heuristics based on feedback."""
    fig, ax = plt.subplots(figsize=(8.6, 4.0), dpi=200)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis("off")
    ctrl = mpatches.FancyBboxPatch((3.5, 5.6), 3.0, 1.4, boxstyle="round,pad=0.05,rounding_size=0.2",
                                    facecolor="#B06A2E", edgecolor="white", linewidth=1.6)
    ax.add_patch(ctrl)
    ax.text(5.0, 6.3, "Hyper-Heuristic\nController (ACO-HH)", ha="center", va="center", fontsize=11,
            color="white", fontweight="bold")

    heuristics = ["H1: 2-opt", "H2: relocate", "H3: swap", "H4: or-opt"]
    for i, h in enumerate(heuristics):
        x = 0.6 + i * 2.35
        box = mpatches.FancyBboxPatch((x, 2.6), 2.0, 1.1, boxstyle="round,pad=0.04,rounding_size=0.12",
                                       facecolor="#2E74B5", edgecolor="white", linewidth=1.3)
        ax.add_patch(box)
        ax.text(x + 1.0, 3.15, h, ha="center", va="center", fontsize=10, color="white", fontweight="bold")
        ax.annotate("", xy=(x + 1.0, 3.7), xytext=(5.0, 5.6),
                    arrowprops=dict(arrowstyle="-|>", color="#8A9BB0", linewidth=1.6))

    ax.annotate(
        "", xy=(6.8, 6.1), xytext=(6.8, 1.6),
        arrowprops=dict(arrowstyle="-|>", color="#3E8E41", linewidth=1.8, connectionstyle="arc3,rad=0.6"),
    )
    ax.text(8.9, 3.8, "pheromone /\nreward feedback\non each heuristic", ha="center", fontsize=9, color="#3E8E41",
            fontweight="bold")
    ax.text(5.0, 0.9, "Picks which low-level heuristic to apply next — searches\nthe space of heuristics, not the space of solutions",
            ha="center", fontsize=10, fontstyle="italic", color="#5A6A7A")
    ax.set_title("Hyper-Heuristics: Choosing Among Heuristics", fontsize=12, fontweight="bold")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_trajectory_overview_image(out_path: Path) -> Path:
    """A single incumbent solution taking a step-by-step trajectory (SANS, PG-CLNS)."""
    fig, ax = plt.subplots(figsize=(8.6, 4.0), dpi=200)
    rng = np.random.default_rng(7)
    steps = 9
    xs = np.linspace(0.5, 9.5, steps)
    ys = 4 + np.cumsum(rng.normal(0, 0.5, steps))
    ys[-1] = min(ys[:-1]) - 0.4
    ax.plot(xs, ys, "-", color="#8A9BB0", linewidth=1.5, zorder=1)
    ax.scatter(xs[:-1], ys[:-1], s=90, color="#B06A2E", zorder=4, edgecolor="white", linewidth=1.2)
    ax.scatter([xs[-1]], [ys[-1]], s=170, color="#3E8E41", zorder=5, edgecolor="white", linewidth=1.5)
    ax.annotate("incumbent", (xs[0], ys[0]), textcoords="offset points", xytext=(-10, 14), fontsize=9.5,
                fontweight="bold", color="#B06A2E")
    ax.annotate("best found", (xs[-1], ys[-1]), textcoords="offset points", xytext=(0, 14), fontsize=9.5,
                fontweight="bold", color="#3E8E41", ha="center")
    for i in range(steps - 2):
        if ys[i + 1] > ys[i]:
            ax.annotate("", xy=(xs[i + 1], ys[i + 1]), xytext=(xs[i], ys[i]),
                        arrowprops=dict(arrowstyle="-|>", color="#C0392B", linewidth=1.2, alpha=0.7))
    ax.text(5, min(ys) - 1.3, "one solution moves step by step; occasional uphill\nmoves escape local optima (e.g. simulated annealing)",
            ha="center", fontsize=10, fontstyle="italic", color="#5A6A7A")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.set_title("Trajectory-Based", fontsize=13, fontweight="bold")
    ax.axis("off")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def generate_population_overview_image(out_path: Path) -> Path:
    """A whole population of candidate solutions evolving together (HGS, PSOMA)."""
    fig, ax = plt.subplots(figsize=(8.6, 4.0), dpi=200)
    rng = np.random.default_rng(3)
    gens = [1.5, 5.0, 8.5]
    gen_labels = ["Generation t", "Generation t+1", "Generation t+2"]
    spread = [1.6, 1.0, 0.55]
    for gi, (gx, spr) in enumerate(zip(gens, spread, strict=True)):
        ys = 4 + rng.normal(0, spr, 8)
        color = ["#B06A2E", "#2E74B5", "#3E8E41"][gi]
        ax.scatter([gx] * 8 + rng.normal(0, 0.15, 8), ys, s=70, color=color, alpha=0.85, zorder=4,
                   edgecolor="white", linewidth=1.0)
        ax.text(gx, 7.2, gen_labels[gi], ha="center", fontsize=10, fontweight="bold", color=color)
        if gi < len(gens) - 1:
            ax.annotate(
                "", xy=(gens[gi + 1] - 0.6, 4), xytext=(gx + 0.6, 4),
                arrowprops=dict(arrowstyle="-|>", color="#8A9BB0", linewidth=1.8),
            )
    ax.text(5, 0.6, "selection + crossover + mutation narrow the population\ntoward better solutions each generation",
            ha="center", fontsize=10, fontstyle="italic", color="#5A6A7A")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis("off")
    ax.set_title("Population-Based", fontsize=13, fontweight="bold")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


ACCENT_HEX = "#2E74B5"

IMAGES_DIR = Path(__file__).resolve().parent / "images"
LINKS_DIR = Path(__file__).resolve().parent / "links"


def generate_vrpp_illustration_fallback(out_path: Path) -> Path:
    """Native fallback for the VRPP illustration: copy the locally stored reference image
    (see archive/gen/images/vrpp_illustration_source.png) instead of a code-drawn diagram."""
    src = IMAGES_DIR / "vrpp_illustration_source.png"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(src, out_path)
    return out_path


def generate_bpc_phases_image(out_path: Path) -> Path:
    """The price and cut phases of Branch-and-Price-and-Cut: the column-generation
    loop between the restricted master problem and the pricing subproblem (left)
    and a cutting plane separating the fractional LP optimum (right)."""
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(6.9, 8.0), dpi=200)

    # ── Left panel: pricing (column generation) loop ──
    ax1.set_xlim(0, 10)
    ax1.set_ylim(0, 10)
    ax1.axis("off")
    ax1.set_title("Price: column generation", fontsize=13, fontweight="bold")

    def _pbox(ax, x, y, w, h, color, label, fontsize=10):
        box = mpatches.FancyBboxPatch(
            (x - w / 2, y - h / 2), w, h, boxstyle="round,pad=0.06,rounding_size=0.15",
            facecolor=color, edgecolor="white", linewidth=1.4, zorder=2,
        )
        ax.add_patch(box)
        ax.text(x, y, label, ha="center", va="center", fontsize=fontsize, color="white",
                fontweight="bold", zorder=3)

    _pbox(ax1, 5, 8.1, 7.6, 2.2, "#1F2D3D", "Restricted Master Problem\nLP over the routes priced so far\nmin Σ c_r λ_r  ·  cover every bin once")
    _pbox(ax1, 5, 2.1, 7.6, 2.2, "#2E74B5", "Pricing subproblem\nshortest path with duals π as arc prices —\nfind a route with negative reduced cost")
    ax1.annotate("", xy=(2.6, 3.4), xytext=(2.6, 6.8),
                 arrowprops=dict(arrowstyle="-|>", color="#B06A2E", linewidth=2.2))
    ax1.text(2.25, 5.1, "duals π", ha="right", va="center", fontsize=10.5, fontweight="bold", color="#B06A2E")
    ax1.annotate("", xy=(7.4, 6.8), xytext=(7.4, 3.4),
                 arrowprops=dict(arrowstyle="-|>", color="#3E8E41", linewidth=2.2))
    ax1.text(7.75, 5.1, "new column\n(route r)", ha="left", va="center", fontsize=10.5,
             fontweight="bold", color="#3E8E41")
    ax1.text(5, 5.1, "repeat until no route\nprices out (< 0)", ha="center", va="center",
             fontsize=9.5, fontstyle="italic", color="#5A6A7A")

    # ── Right panel: cutting plane on the LP relaxation ──
    ax2.set_xlim(-0.5, 10.5)
    ax2.set_ylim(-0.5, 9.5)
    ax2.axis("off")
    ax2.set_title("Cut: tighten the LP relaxation", fontsize=13, fontweight="bold")
    # integer lattice
    for gx in range(10):
        for gy in range(9):
            ax2.plot([gx], [gy], marker="o", markersize=2.4, color="#C4CDD6", zorder=1)
    # LP feasible polytope
    poly = np.array([[0.5, 0.6], [8.9, 1.1], [9.4, 4.4], [6.4, 8.3], [1.4, 7.4]])
    ax2.add_patch(mpatches.Polygon(poly, closed=True, facecolor="#DCE7F5", edgecolor="#2E74B5",
                                    linewidth=2.0, alpha=0.9, zorder=2))
    # integer points inside the polytope
    from matplotlib.path import Path as _MplPath
    ppath = _MplPath(poly)
    for gx in range(10):
        for gy in range(9):
            if ppath.contains_point((gx, gy), radius=-1e-9):
                ax2.plot([gx], [gy], marker="o", markersize=4.2, color="#1F2D3D", zorder=3)
    # fractional LP optimum near a vertex
    frac = (8.05, 6.15)
    ax2.plot([frac[0]], [frac[1]], marker="*", markersize=17, color="#C0392B", zorder=5)
    ax2.annotate("fractional LP\noptimum x*", frac, textcoords="offset points", xytext=(14, 12),
                 fontsize=10, fontweight="bold", color="#C0392B")
    # cutting plane separating x* from the integer hull
    ax2.plot([4.4, 10.2], [8.9, 2.5], color="#B06A2E", linewidth=2.6, linestyle="--", zorder=4)
    ax2.text(-0.3, 8.9, "valid inequality cuts off x*\n(all integer points kept)",
             fontsize=10, fontweight="bold", color="#B06A2E", ha="left", va="center")

    fig.suptitle("The price and cut phases at every tree node", fontsize=14, fontweight="bold")
    fig.tight_layout(rect=(0, 0, 1, 0.95))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


NATIVE_DIAGRAM_BUILDERS = {
    "bb_tree.png": generate_bb_tree_image,
    "bpc_phases.png": generate_bpc_phases_image,
    "ls_operators.png": generate_ls_operators_image,
    "knapsack_illustration.png": generate_knapsack_image,
    "lookahead_flow.png": generate_lookahead_flow_image,
    "framework_objective.png": generate_framework_objective_image,
    "metaheuristic_overview.png": generate_metaheuristic_overview_image,
    "hyperheuristic_overview.png": generate_hyperheuristic_overview_image,
    "trajectory_overview.png": generate_trajectory_overview_image,
    "population_overview.png": generate_population_overview_image,
    "vrpp_illustration.png": generate_vrpp_illustration_fallback,
}


def _load_reference_links() -> dict:
    """Parse archive/gen/links/reference_image_links.xml -> {key: {"url": ...}}."""
    path = LINKS_DIR / "reference_image_links.xml"
    if not path.exists():
        return {}
    root = ET.parse(path).getroot()
    links: dict = {}
    for img in root.findall("image"):
        key = img.get("key")
        url_el = img.find("url")
        if key and url_el is not None and url_el.text:
            links[key] = {"url": url_el.text.strip()}
    return links


def ensure_reference_images(figures_dir: Path, image_mode: str) -> None:
    """Populate the native (or fetched) reference-style diagrams used across the deck."""
    links = _load_reference_links() if image_mode == "fetch" else {}
    for fname, builder in NATIVE_DIAGRAM_BUILDERS.items():
        out_path = figures_dir / fname
        if image_mode == "fetch":
            key = Path(fname).stem
            url = links.get(key, {}).get("url")
            if url:
                try:
                    import urllib.request

                    out_path.parent.mkdir(parents=True, exist_ok=True)
                    urllib.request.urlretrieve(url, out_path)  # noqa: S310
                    print(f"  Fetched: {out_path} <- {url}")
                    continue
                except Exception as exc:  # noqa: BLE001
                    print(f"  [WARN] Fetch failed for {key} ({exc}); falling back to native diagram")
        builder(out_path)
        print(f"  Saved (native): {out_path}")


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument(
        "--figures-dir",
        default="public/figures/simulation/30d",
        help="Directory with the simulation analysis figures to embed",
    )
    p.add_argument("--out", default="assets/windows/wsmart_route_results.pptx", help="Destination .pptx path")
    p.add_argument("--author", default=None, help="Presenting author shown on the cover slide")
    p.add_argument(
        "--coauthors", default=None, help="Optional semicolon-separated co-authors (shown with less emphasis)"
    )
    p.add_argument("--groups", default=None, help="Optional semicolon-separated research groups of the authors")
    p.add_argument(
        "--results-table",
        default="30d",
        choices=["30d", "90d", "all", "none"],
        help="Full results table slide: a single horizon, 'all' horizons "
        "(horizon added to the column hierarchy), or 'none' to omit the slide",
    )
    p.add_argument(
        "--results-table-split",
        default="none",
        choices=["none", "strategy", "constructor", "improver"],
        help="Partition the results table by this column-hierarchy level: it becomes the "
        "outermost grouping and each value gets its own stacked partial table on the slide "
        "(e.g. 'improver' -> one partial table for CLS, one for FTSP)",
    )
    p.add_argument(
        "--speaker-script",
        action="store_true",
        help="Also generate a per-slide speaker script .docx alongside the .pptx",
    )
    p.add_argument(
        "--speaker-script-out",
        default=None,
        help="Destination .docx path (default: same name as --out, .docx extension)",
    )
    p.add_argument(
        "--image-mode",
        default="native",
        choices=["native", "fetch"],
        help="How to source the conceptual diagrams (B&B tree, GA cycle, 2-opt swap, VRPP illustration): "
        "'native' draws them in-house / copies the locally stored reference image (default, no "
        "licensing/network concerns); 'fetch' downloads the images listed in "
        "links/reference_image_links.xml instead.",
    )
    p.add_argument("--excel", action="store_true", help="Also export the CLS results table as an Excel workbook")
    return p.parse_args()


def export_results_excel(out_path: Path, results_table: str = "30d") -> None:
    """Export the CLS-only results matrix to an Excel workbook (item 9 Excel template)."""
    import openpyxl
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    results_config = load_json("simulation_analysis_config.json")
    horizon_specs = sorted(results_config["horizons"], key=lambda h: h["days"])

    if results_table == "all":
        all_data: list[tuple] = []
        for spec in horizon_specs:
            csv_path = Path(spec["csv"])
            if not csv_path.exists():
                continue
            df = filter_data(load_horizon_csv(csv_path), results_config)
            if df.empty:
                continue
            theme = {"name": "light", "fg": "#1F2D3D"}
            ctx = build_context(df, results_config, theme, spec["days"])
            df_cls = df[df["improver"].str.upper() == "CLS"]
            if df_cls.empty:
                df_cls = df
            dfm_cls = aggregate(df_cls)
            ctx_cls = {**ctx, "improvers": list(df_cls["improver"].unique())}
            rk, ck, cs = build_full_results_matrix(dfm_cls, ctx_cls, horizon_label=f"{spec['days']}d")
            all_data.append((spec["days"], rk, ck, cs, ctx_cls))
        if not all_data:
            print("  [WARN] No data for Excel export")
            return
        _, row_keys, col_keys, cells, ctx_ref = all_data[0]
        for _, rk, ck, cs, _ in all_data[1:]:
            row_keys = list({*row_keys, *rk})
            col_keys = col_keys + ck
            cells.update(cs)
        row_keys.sort()
    else:
        days = int(results_table.rstrip("d"))
        spec = next((s for s in horizon_specs if s["days"] == days), None)
        if spec is None:
            return
        csv_path = Path(spec["csv"])
        if not csv_path.exists():
            print(f"  [WARN] CSV not found: {csv_path} — skipping Excel export")
            return
        df = filter_data(load_horizon_csv(csv_path), results_config)
        if df.empty:
            return
        theme = {"name": "light", "fg": "#1F2D3D"}
        ctx_ref = build_context(df, results_config, theme, days)
        df_cls = df[df["improver"].str.upper() == "CLS"]
        if df_cls.empty:
            df_cls = df
        dfm_cls = aggregate(df_cls)
        ctx_cls = {**ctx_ref, "improvers": list(df_cls["improver"].unique())}
        row_keys, col_keys, cells = build_full_results_matrix(dfm_cls, ctx_cls)

    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Results"

    HDR_FILL = PatternFill("solid", fgColor="1F2D3D")
    HDR_FONT = Font(color="FFFFFF", bold=True)
    ALT_FILL = PatternFill("solid", fgColor="EEF2F7")
    BEST_OV_FILL = PatternFill("solid", fgColor="C6EFCE")
    BEST_KG_FILL = PatternFill("solid", fgColor="BDD7EE")
    CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)

    n_row_labels = 3  # Region, N, Dist
    header_row = 1
    data_row_start = header_row + 1

    # Write column headers
    for ci, ck in enumerate(col_keys, start=n_row_labels + 1):
        cell = ws.cell(row=header_row, column=ci, value=" / ".join(str(x) for x in ck))
        cell.fill = HDR_FILL
        cell.font = HDR_FONT
        cell.alignment = CENTER

    # Row label headers
    for li, label in enumerate(["Region", "N", "Distribution"], start=1):
        cell = ws.cell(row=header_row, column=li, value=label)
        cell.fill = HDR_FILL
        cell.font = HDR_FONT
        cell.alignment = CENTER

    # Write data rows
    global_best = compute_global_best(row_keys, col_keys, cells)
    prev: tuple = ()
    for ri, rk in enumerate(row_keys, start=data_row_start):
        fill = ALT_FILL if ri % 2 == 0 else PatternFill()
        for li, val in enumerate(rk, start=1):
            display = val if prev[:li] != rk[:li] else ""
            cell = ws.cell(row=ri, column=li, value=display)
            cell.alignment = CENTER
            if fill:
                cell.fill = fill
        prev = rk
        best = global_best.get(rk, {})
        for ci, ck in enumerate(col_keys, start=n_row_labels + 1):
            raw = cells.get((rk, ck), "—")
            clean = raw.replace("<br>", "\n")
            cell = ws.cell(row=ri, column=ci, value=clean)
            cell.alignment = CENTER
            if fill:
                cell.fill = fill
            if ck == best.get("ov"):
                cell.fill = BEST_OV_FILL
            elif ck == best.get("kg"):
                cell.fill = BEST_KG_FILL

    # Column widths
    for ci in range(1, n_row_labels + 1):
        ws.column_dimensions[get_column_letter(ci)].width = 14
    for ci in range(n_row_labels + 1, len(col_keys) + n_row_labels + 1):
        ws.column_dimensions[get_column_letter(ci)].width = 18
    ws.row_dimensions[header_row].height = 40

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    print(f"  Excel template written: {out_path}")


def main() -> None:
    args = parse_args()
    content = load_json("presentation_content.json")
    figures_dir = Path(args.figures_dir)
    if not figures_dir.is_dir():
        raise SystemExit(f"Figures dir not found: {figures_dir} — run gen_simulation_analysis.py first")

    # Generate the QA route illustration and the conceptual diagrams (incl. the VRPP illustration) dynamically
    generate_qa_route_image(figures_dir / "qa_route_illustration.png")
    ensure_reference_images(figures_dir, args.image_mode)

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    coauthors = [c.strip() for c in args.coauthors.split(";")] if args.coauthors else None
    groups = [g.strip() for g in args.groups.split(";")] if args.groups else None
    builder = DeckBuilder(
        content, figures_dir, args.author, coauthors, groups, args.results_table, args.results_table_split
    )
    prs = builder.build()
    prs.save(str(out))
    print(f"Written: {out} ({len(prs.slides._sldIdLst)} slides)")

    if args.speaker_script:
        script_out = Path(args.speaker_script_out) if args.speaker_script_out else out.with_suffix(".docx")
        gen_speaker_script(content["title"], builder.author, builder.slide_scripts, script_out)
        appendix_md = Path(__file__).resolve().parent / "appendix_notes.md"
        if appendix_md.exists():
            append_markdown_appendix(script_out, appendix_md)
            print(f"Written: {script_out} ({len(builder.slide_scripts)} slide scripts + Q&A appendix)")
        else:
            print(f"Written: {script_out} ({len(builder.slide_scripts)} slide scripts)")

    if getattr(args, "excel", False):
        excel_out = out.with_suffix(".xlsx")
        export_results_excel(excel_out, results_table=args.results_table)
        print(f"Written: {excel_out}")


if __name__ == "__main__":
    main()
