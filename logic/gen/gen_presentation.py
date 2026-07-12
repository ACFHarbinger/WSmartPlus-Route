"""
Generate the WSmart+ Route results PowerPoint presentation.

Builds a 20-slide deck under assets/windows/ following the agreed structure:

  1.  Cover (title/authors/affiliations, as in the conference abstract)
  2.  Index / agenda (condensed)
  3.  The VRPP problem                          (equation + caption)
  4.  Objective of this work                    (bullets)
  5.  Routing simulator philosophy              (pipeline diagram + caption)
  6.  Mandatory node selection strategies       (equation + caption)
  7.  Big picture: policy configuration space   (grid diagram + caption)
  8.  Exact methods (BPC + SWC-TCF)             (equation + figure + caption)
  9.  Meta-heuristics & hyper-heuristics        (equation + caption, all algorithms)
  10. CLS vs Fast-TSP route improvers           (equation + figure + caption)
  11. Design of experiments                     (tree diagram + caption)
  12. Pareto front plot (log X, per-scenario coloured fronts, figure caption)
  13. Summary KPI plots (stacked vertically, full width, figure caption)
  14. Strategy trade-off bubble chart (figure caption)
  15. Per-scenario heatmaps (30 days, figure caption)
  16. Overflow + kg/km policy × scenario heatmaps (90 days, side legend + caption)
  17. Route improver bubble chart (figure caption)
  18. Full results table (user-selected horizon, or all horizons — table caption)
  19. Conclusions, limitations & future work    (radar figure + caption)
  20. End / Q&A (figure)

Slide text, equations (LaTeX, embedded as native editable OMML equations)
and captions live in json/presentation_content.json; result figures are
pulled from the simulation analysis figures directories (30-day set by
default; individual figure slides may override via "figures_dir"). The full
results table (slide 18) is built directly from the horizon CSV(s)
referenced in json/simulation_analysis_config.json.

A per-slide speaker script can also be generated as a .docx (see
gen_speaker_script / --speaker-script), rendered via docxtpl from a template
under logic/gen/templates/.

Usage
-----
    uv run python logic/gen/gen_presentation.py
    uv run python logic/gen/gen_presentation.py \\
        --figures-dir public/figures/simulation/30d \\
        --out assets/windows/wsmart_route_results.pptx \\
        --author "Afonso Fernandes" \\
        --coauthors "Jane Doe;John Smith" \\
        --groups "ISR Coimbra;INESC-ID Lisboa" \\
        --results-table 30d \\
        --speaker-script
"""

from __future__ import annotations

import argparse
import re
import subprocess
import tempfile
import textwrap
import xml.sax.saxutils as saxutils
import zipfile
from pathlib import Path

import matplotlib
import numpy as np

matplotlib.use("Agg")
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
from gen_simulation_analysis import (  # pyrefly: ignore [missing-import]
    _group_spans,
    aggregate,
    build_context,
    build_full_results_matrix,
    filter_data,
    load_horizon_csv,
)
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationClass
from pptx.util import Emu, Inches, Pt
from report_utils import load_json  # pyrefly: ignore [missing-import]

# ── Editable (OMML) equations ───────────────────────────────────────────────
# Equations are converted from LaTeX to Office Math Markup (OMML) via pandoc
# and embedded directly in the slide XML (mc:AlternateContent / a14:m), so
# they open as native, editable equation objects in PowerPoint rather than
# rasterised images. Non-Microsoft viewers fall back to a plain-text run.
_MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _latex_to_omath(latex: str) -> etree._Element:
    """Convert one LaTeX equation line to an <m:oMath> element via pandoc."""
    with tempfile.TemporaryDirectory() as td:
        md_path, docx_path = Path(td) / "eq.md", Path(td) / "eq.docx"
        md_path.write_text(f"$${latex}$$", encoding="utf-8")
        subprocess.run(["pandoc", str(md_path), "-o", str(docx_path)], check=True, capture_output=True)
        with zipfile.ZipFile(docx_path) as z:
            doc_xml = z.read("word/document.xml")
    root = etree.fromstring(doc_xml)
    omath = root.find(f".//{{{_MATH_NS}}}oMath")
    # Drop WordprocessingML control-run props (word/pptx sizing units differ);
    # the equation instead inherits the paragraph's own default run formatting.
    for ctrl_pr in omath.findall(f".//{{{_MATH_NS}}}ctrlPr"):
        for child in list(ctrl_pr):
            ctrl_pr.remove(child)
    return omath


def _plain_fallback(latex: str) -> str:
    """A readable, non-mathematical fallback string for non-OOXML-math viewers."""
    text = re.sub(r"\\(mathbf|mathrm|textbf|text|mathcal|left|right)\b", "", latex)
    text = re.sub(r"[\\{}$]", "", text).replace("\\,", " ").replace("\\ ", " ").replace("\\quad", "  ")
    return re.sub(r"\s+", " ", text).strip()


def _apply_equation_to_paragraph(p, latex: str, size_pt: int = 22, color: str = "1F2D3D", align: str = "ctr"):
    """Turn an (empty) paragraph into a native, editable PowerPoint equation."""
    omath = _latex_to_omath(latex)
    p_el = p._p
    pPr = etree.SubElement(p_el, qn("a:pPr"))
    pPr.set("algn", align)
    def_rpr = etree.SubElement(pPr, qn("a:defRPr"))
    def_rpr.set("sz", str(size_pt * 100))
    fill = etree.SubElement(def_rpr, qn("a:solidFill"))
    etree.SubElement(fill, qn("a:srgbClr")).set("val", color)
    alt_xml = (
        f'<mc:AlternateContent xmlns:mc="{_MC_NS}" xmlns:a14="{_A14_NS}" xmlns:m="{_MATH_NS}" xmlns:a="{_A_NS}">'
        f'<mc:Choice xmlns:a14="{_A14_NS}" Requires="a14">'
        f'<a14:m><m:oMathPara xmlns:m="{_MATH_NS}">{etree.tostring(omath, encoding="unicode")}</m:oMathPara></a14:m>'
        f"</mc:Choice>"
        f'<mc:Fallback><a:r><a:rPr lang="en-US" sz="{size_pt * 100}">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:rPr>'
        f"<a:t>{saxutils.escape(_plain_fallback(latex))}</a:t></a:r></mc:Fallback>"
        f"</mc:AlternateContent>"
    )
    p_el.append(etree.fromstring(alt_xml))
    return p


def add_equation_paragraphs(text_frame, lines: list[str], size_pt: int = 22, color: str = "1F2D3D", align: str = "ctr"):
    """Render each line of `lines` as its own native-equation paragraph in `text_frame`."""
    for i, latex in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        _apply_equation_to_paragraph(p, latex, size_pt=size_pt, color=color, align=align)


# 16:9 slide geometry
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT = RGBColor(0x2E, 0x74, 0xB5)
DARK = RGBColor(0x1F, 0x2D, 0x3D)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
LIGHT_TXT = RGBColor(0xC9, 0xD6, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PIPELINE_STAGES = [
    ("Mandatory bin\nselection strategy", "LA · LM · SL\nwhich bins today?"),
    ("Route\nconstructor", "exact · meta-h. · hyper-h.\nbuild the routes"),
    ("Acceptance criterion\n(optional)", "constructor-dependent\ne.g. BMC, OI"),
    ("Route improver\n(optional)", "CLS · FTSP\npost-optimisation"),
]


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _wrap_label(label, width: int) -> str:
    """Wrap a table header/cell label onto multiple lines so it fits its cell."""
    return "\n".join(textwrap.wrap(str(label), width=width)) or str(label)


_CELL_RE = re.compile(r"([-\d.]+(?:±[-\d.]+)?)\s*ov\s*(?:<br>|\n)\s*([-\d.]+(?:±[-\d.]+)?)\s*kg/km")


def _parse_result_cell(text: str):
    """Split a formatted result cell ('X ov<br>Y kg/km') into (overflow, kgkm) as (raw_str, numeric_mean)."""
    m = _CELL_RE.match(text)
    if not m:
        return None, None

    def _pair(raw: str):
        try:
            return raw, float(raw.split("±")[0])
        except ValueError:
            return raw, None

    return _pair(m.group(1)), _pair(m.group(2))


def render_hier_table_image(
    row_keys: list[tuple],
    col_keys: list[tuple],
    cells: dict[tuple, str],
    row_labels: list[str],
    out_path: Path,
    col_lookup_keys: list[tuple] | None = None,
    partition_label: str | None = None,
) -> Path:
    """
    Render a hierarchical (merged-header) results table as a PNG, in the style
    of a LaTeX multi-row/multi-column table: row groups (graph size > region >
    distribution) and column groups (strategy > constructor > improver, plus
    horizon if present) are drawn as merged header spans.

    `col_lookup_keys`, if given, is parallel to `col_keys` and used instead of
    it to look cells up in `cells` — lets a hierarchy level be dropped from the
    displayed header (e.g. when the table has been partitioned by that level)
    while the underlying data lookup still uses the full key.
    """
    col_lookup_keys = col_lookup_keys if col_lookup_keys is not None else col_keys
    n_row_levels = len(row_labels)
    n_col_levels = len(col_keys[0])
    n_rows = len(row_keys)
    n_cols = len(col_keys)

    # Reserve an extra banner row at the top when a partition label is supplied
    banner_h = 0.5 if partition_label else 0.0

    cell_w, cell_h, label_w, header_h = 1.05, 0.5, 1.3, 0.5
    fontsize = max(5.5, min(9, 260 / max(n_cols, 1)))

    x0 = label_w * n_row_levels
    y0 = banner_h + header_h * n_col_levels
    total_h = banner_h + header_h * n_col_levels + cell_h * n_rows
    total_w = x0 + cell_w * n_cols
    fig_w = min(total_w, 42)
    fig_h = min(total_h, 32)

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.set_xlim(0, total_w)
    ax.set_ylim(0, total_h)
    ax.invert_yaxis()
    ax.axis("off")

    # Banner row spanning all columns (including the row-label area)
    if partition_label:
        ax.add_patch(
            mpatches.Rectangle(
                (0, 0),
                total_w,
                banner_h,
                facecolor="#0D1B2A",
                edgecolor="none",
            )
        )
        ax.text(
            total_w / 2,
            banner_h / 2,
            partition_label,
            ha="center",
            va="center",
            fontsize=fontsize + 1,
            color="white",
            fontweight="bold",
        )

    header_colors = ["#1F2D3D", "#2E74B5", "#5A6A7A", "#8A9BB0"]
    for lvl in range(n_col_levels):
        y_top = banner_h + lvl * header_h
        for start, end, label in _group_spans(col_keys, lvl):
            xs, xe = x0 + start * cell_w, x0 + end * cell_w
            ax.add_patch(
                mpatches.Rectangle(
                    (xs, y_top),
                    xe - xs,
                    header_h,
                    facecolor=header_colors[lvl % len(header_colors)],
                    edgecolor="white",
                    linewidth=0.8,
                )
            )
            wrap_chars = max(6, int((xe - xs) / cell_w) * 10)
            ax.text(
                (xs + xe) / 2,
                y_top + header_h / 2,
                _wrap_label(label, wrap_chars),
                ha="center",
                va="center",
                fontsize=fontsize,
                color="white",
                fontweight="bold",
            )

    for lvl in range(n_row_levels):
        x_left = lvl * label_w
        for start, end, label in _group_spans(row_keys, lvl):
            ys, ye = y0 + start * cell_h, y0 + end * cell_h
            ax.add_patch(
                mpatches.Rectangle(
                    (x_left, ys),
                    label_w,
                    ye - ys,
                    facecolor="#F0F4FA",
                    edgecolor="#5A6A7A",
                    linewidth=0.6,
                )
            )
            ax.text(
                x_left + label_w / 2,
                (ys + ye) / 2,
                _wrap_label(label, 10),
                ha="center",
                va="center",
                fontsize=fontsize,
                color="#1F2D3D",
                fontweight="bold",
            )

    for ri, rk in enumerate(row_keys):
        parsed = {ci: _parse_result_cell(cells.get((rk, lk), "—")) for ci, lk in enumerate(col_lookup_keys)}
        ov_vals = {ci: v[0][1] for ci, v in parsed.items() if v[0] and v[0][1] is not None}
        kg_vals = {ci: v[1][1] for ci, v in parsed.items() if v[1] and v[1][1] is not None}
        best_ov_ci = min(ov_vals, key=ov_vals.get) if ov_vals else None  # pyrefly: ignore [no-matching-overload]
        best_kg_ci = max(kg_vals, key=kg_vals.get) if kg_vals else None  # pyrefly: ignore [no-matching-overload]
        for ci in range(n_cols):
            xs, ys = x0 + ci * cell_w, y0 + ri * cell_h
            ax.add_patch(
                mpatches.Rectangle(
                    (xs, ys),
                    cell_w,
                    cell_h,
                    fill=False,
                    edgecolor="#CCCCCC",
                    linewidth=0.4,
                )
            )
            ov, kg = parsed.get(ci, (None, None))
            if ov is None and kg is None:
                ax.text(xs + cell_w / 2, ys + cell_h / 2, "—", ha="center", va="center", fontsize=fontsize * 0.85)
                continue
            ov_color, ov_weight = ("#1A7A34", "bold") if ci == best_ov_ci else ("#333333", "normal")
            kg_color, kg_weight = ("#1A7A34", "bold") if ci == best_kg_ci else ("#333333", "normal")
            ax.text(
                xs + cell_w / 2,
                ys + cell_h * 0.3,
                ov[0] if ov else "—",
                ha="center",
                va="center",
                fontsize=fontsize * 0.85,
                color=ov_color,
                fontweight=ov_weight,
            )
            ax.text(
                xs + cell_w / 2,
                ys + cell_h * 0.7,
                kg[0] if kg else "—",
                ha="center",
                va="center",
                fontsize=fontsize * 0.85,
                color=kg_color,
                fontweight=kg_weight,
            )

    fig.savefig(out_path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


TEMPLATES_DIR = Path(__file__).resolve().parent / "templates"


def gen_speaker_script(deck_title: str, author: str, slides: list[dict], out_path: Path) -> Path:
    """Render a per-slide speaker script .docx from the docxtpl template."""
    from docxtpl import DocxTemplate

    tpl = DocxTemplate(str(TEMPLATES_DIR / "speaker_script.docx"))
    tpl.render({"deck_title": deck_title, "author": author, "slides": slides})
    out_path.parent.mkdir(parents=True, exist_ok=True)
    tpl.save(str(out_path))
    return out_path


class DeckBuilder:
    def __init__(
        self,
        content: dict,
        figures_dir: Path,
        author: str | None,
        coauthors: list[str] | None = None,
        groups: list[str] | None = None,
        results_table: str = "30d",
        results_table_split: str = "none",
    ):
        self.content = content
        self.figures_dir = figures_dir
        self.author = author or content["author"]
        self.coauthors = coauthors if coauthors is not None else content.get("coauthors", [])
        self.groups = groups if groups is not None else content.get("research_groups", [])
        self.results_table = results_table
        self.results_table_split = results_table_split
        self._results_config: dict = {}
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]
        self._tmp = Path(tempfile.mkdtemp(prefix="wsr_pptx_"))
        self._eq_count = 0
        self._fig_count = 0
        self._tab_count = 0
        self.slide_scripts: list[dict] = []

    def _record_script(self, title: str, paragraphs: list[str]) -> None:
        """Append a speaker-script entry for the slide just built."""
        self.slide_scripts.append(
            {
                "number": len(self.slide_scripts) + 1,
                "title": title,
                "script": "\n\n".join(p for p in paragraphs if p),
            }
        )

    # ── Building blocks ─────────────────────────────────────────────────────────

    def _new_slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _title_bar(self, slide, title: str) -> None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))
        _fill(bar, DARK)
        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE

    def _bullets(self, slide, bullets: list[str], left=None, top=None, width=None, height=None, size=16) -> None:
        left = left if left is not None else Inches(0.6)
        top = top if top is not None else Inches(1.3)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        _, tf = _textbox(slide, left, top, width, height or (SLIDE_H - top - Inches(0.4)))
        for i, text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {text}"
            p.font.size = Pt(size)
            p.font.color.rgb = DARK
            p.space_after = Pt(10)

    def _picture_fit(self, slide, path: Path, left, top, max_w, max_h) -> None:
        from PIL import Image

        with Image.open(path) as im:
            w_px, h_px = im.size
        ratio = min(max_w / w_px, max_h / h_px)
        w, h = int(w_px * ratio), int(h_px * ratio)
        slide.shapes.add_picture(str(path), left + Emu(int((max_w - w) / 2)), top + Emu(int((max_h - h) / 2)), w, h)

    def _caption_box(self, slide, label: str, text: str, top=None, left=None, width=None, height=None) -> None:
        """A numbered caption ('**Figure N:** ...' / '**Equation N:** ...' / '**Table N:** ...')."""
        top = top if top is not None else SLIDE_H - Inches(0.55)
        left = left if left is not None else Inches(0.6)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        box_h = height if height is not None else Inches(0.9)
        _, tf = _textbox(slide, left, top, width, box_h)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{label}: "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = DARK
        r2 = p.add_run()
        r2.text = text
        r2.font.italic = True
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED

    def _figure_caption(self, slide, text: str, top=None, left=None, width=None, height=None) -> None:
        self._fig_count += 1
        self._caption_box(slide, f"Figure {self._fig_count}", text, top=top, left=left, width=width, height=height)

    def _equation_caption(self, slide, text: str, top=None, left=None, width=None) -> None:
        self._caption_box(slide, f"Equation {self._eq_count}", text, top=top, left=left, width=width)

    def _table_caption(self, slide, text: str, top=None) -> None:
        self._tab_count += 1
        self._caption_box(slide, f"Table {self._tab_count}", text, top=top)

    def _equation_focus(self, slide, lines: list[str], left=None, top=None, width=None, size_pt: int = 22):
        """Place native, editable equation lines as the slide's visual focus."""
        self._eq_count += 1
        left = left if left is not None else Inches(0.6)
        top = top if top is not None else Inches(1.25)
        width = width if width is not None else SLIDE_W - Inches(1.2)
        area_h = Inches(0.62 * len(lines) + 0.4)
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, area_h)
        _fill(band, RGBColor(0xF0, 0xF4, 0xFA))
        band.shadow.inherit = False
        box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), area_h - Inches(0.2))
        tf = box.text_frame
        tf.word_wrap = True
        add_equation_paragraphs(tf, lines, size_pt=size_pt)
        return top + area_h + Inches(0.2)

    def _pipeline_diagram(self, slide) -> None:
        """Draw the policy pipeline as chevron stages with captions."""
        n = len(PIPELINE_STAGES)
        gap = Inches(0.25)
        left0 = Inches(0.6)
        total_w = SLIDE_W - Inches(1.2)
        stage_w = int((total_w - gap * (n - 1)) / n)
        top = Inches(1.7)
        h = Inches(1.5)
        for i, (name, caption) in enumerate(PIPELINE_STAGES):
            left = left0 + i * (stage_w + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, stage_w, h)
            optional = "optional" in name
            _fill(shape, ACCENT if not optional else MUTED)
            shape.shadow.inherit = False
            tf = shape.text_frame
            tf.word_wrap = True
            for li, line in enumerate(name.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            _, ctf = _textbox(slide, left, top + h + Inches(0.1), stage_w, Inches(0.9))
            for li, line in enumerate(caption.split("\n")):
                p = ctf.paragraphs[0] if li == 0 else ctf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
                p.font.color.rgb = MUTED
                p.alignment = PP_ALIGN.CENTER
        sim = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.1), top + h + Inches(1.2), SLIDE_W - Inches(6.2), Inches(0.8)
        )
        _fill(sim, DARK)
        sim.shadow.inherit = False
        p = sim.text_frame.paragraphs[0]
        p.text = "Multi-day simulator: region × graph size × distribution"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    def _policy_grid_diagram(self, slide) -> None:
        """Draw the Mandatory Selection x Route Constructor x Route Improver big picture."""
        columns = [
            (
                "Mandatory Selection",
                [
                    "Look-Ahead (LA)",
                    "Last-Minute (LM, CF70)",
                    "Last-Minute (LM, CF90)",
                    "Service-Level (SL1)",
                    "Service-Level (SL2)",
                ],
                "flat",
            ),
            (
                "Route Constructor",
                [
                    ("Exact Methods", ["BPC", "SWC-TCF"]),
                    ("Meta-Heuristics", ["HGS", "ALNS", "SANS", "PG-CLNS"]),
                    ("Hyper-Heuristics", ["PSOMA", "ACO-HH"]),
                ],
                "grouped",
            ),
            ("Route Improver", ["Fast-TSP", "Local Search (CLS)"], "flat"),
        ]
        top, bottom = Inches(1.15), Inches(5.75)
        gap = Inches(0.3)
        left0 = Inches(0.6)
        col_w = int((SLIDE_W - Inches(1.2) - gap * (len(columns) - 1)) / len(columns))
        header_h = Inches(0.55)
        for ci, (header, items, kind) in enumerate(columns):
            left = left0 + ci * (col_w + gap)
            hd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, header_h)
            _fill(hd, DARK)
            hd.shadow.inherit = False
            p = hd.text_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            body_top = top + header_h + Inches(0.15)
            body_h = bottom - body_top
            if kind == "flat":
                item_gap = Inches(0.12)
                item_h = int((body_h - item_gap * (len(items) - 1)) / len(items))
                for ii, label in enumerate(items):
                    it_top = body_top + ii * (item_h + item_gap)
                    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, it_top, col_w, item_h)
                    _fill(box, ACCENT)
                    box.shadow.inherit = False
                    p = box.text_frame.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(13)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
            else:
                grp_gap = Inches(0.15)
                grp_h = int((body_h - grp_gap * (len(items) - 1)) / len(items))
                for gi, (grp_label, sub_items) in enumerate(items):
                    grp_top = body_top + gi * (grp_h + grp_gap)
                    grp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, grp_top, col_w, grp_h)
                    grp.fill.background()
                    grp.line.color.rgb = MUTED
                    grp.line.width = Pt(1)
                    grp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                    grp.shadow.inherit = False
                    tf = grp.text_frame
                    tf.word_wrap = True
                    tf.margin_top = Inches(0.05)
                    p = tf.paragraphs[0]
                    p.text = grp_label
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = ACCENT
                    p.alignment = PP_ALIGN.CENTER
                    pi = tf.add_paragraph()
                    pi.text = "  ·  ".join(sub_items)
                    pi.font.size = Pt(11)
                    pi.font.color.rgb = DARK
                    pi.alignment = PP_ALIGN.CENTER

    def _doe_tree_diagram(self, slide) -> None:
        """Draw the design-of-experiments tree: horizon -> scenario -> distribution."""
        horizons = [
            (
                "30 Days",
                [
                    ("RM-100", ["Empirical", "Gamma-3"]),
                    ("RM-170", ["Empirical", "Gamma-3"]),
                    ("FFZ-350", ["Empirical", "Gamma-3"]),
                ],
            ),
            (
                "90 Days\n(Pareto-front policies only)",
                [
                    ("RM-100", ["Empirical", "Gamma-3"]),
                    ("RM-170", ["Empirical", "Gamma-3"]),
                    ("FFZ-350", ["Empirical", "Gamma-3"]),
                ],
            ),
        ]
        top = Inches(1.25)
        root_w, root_h = Inches(2.2), Inches(0.5)
        root_left = (SLIDE_W - root_w) / 2
        root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, root_left, top, root_w, root_h)
        _fill(root, DARK)
        root.shadow.inherit = False
        p = root.text_frame.paragraphs[0]
        p.text = "Simulation Runs"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        n_h = len(horizons)
        gap_h = Inches(0.4)
        hz_top = top + root_h + Inches(0.35)
        hz_w = int((SLIDE_W - Inches(1.2) - gap_h * (n_h - 1)) / n_h)
        hz_h = Inches(0.6)
        left0 = Inches(0.6)
        for hi, (hz_label, scenarios) in enumerate(horizons):
            hz_left = left0 + hi * (hz_w + gap_h)
            hz_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, hz_left, hz_top, hz_w, hz_h)
            _fill(hz_box, ACCENT)
            hz_box.shadow.inherit = False
            tf = hz_box.text_frame
            tf.word_wrap = True
            for li, line in enumerate(hz_label.split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

            n_s = len(scenarios)
            gap_s = Inches(0.15)
            sc_top = hz_top + hz_h + Inches(0.3)
            sc_w = int((hz_w - gap_s * (n_s - 1)) / n_s)
            sc_h = Inches(0.5)
            for si, (sc_label, dists) in enumerate(scenarios):
                sc_left = hz_left + si * (sc_w + gap_s)
                sc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sc_left, sc_top, sc_w, sc_h)
                _fill(sc_box, MUTED)
                sc_box.shadow.inherit = False
                p = sc_box.text_frame.paragraphs[0]
                p.text = sc_label
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

                dist_top = sc_top + sc_h + Inches(0.1)
                dist_h = Inches(0.4)
                _, dtf = _textbox(slide, sc_left, dist_top, sc_w, dist_h)
                p = dtf.paragraphs[0]
                p.text = " / ".join(dists)
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER

    def _figure_slide(self, spec: dict) -> None:
        slide = self._new_slide()
        self._title_bar(slide, spec["title"])
        fig_dir = Path(spec["figures_dir"]) if spec.get("figures_dir") else self.figures_dir
        figures = spec.get("figures") or [spec.get("figure")]
        resolved = []
        for name in figures:
            p = fig_dir / name
            if not p.exists() and spec.get("fallback_figure"):
                p = fig_dir / spec["fallback_figure"]
            if p.exists():
                resolved.append(p)
            else:
                print(f"  [WARN] Figure not found: {p}")
        area_top = Inches(1.15)
        bottom_legend = spec.get("bottom_legend", False)
        caption_top = None
        caption_h = None
        legend_top = None
        legend_h = None
        area_h = None
        if bottom_legend:
            fig_area_h = Inches(4.2)
            caption_top = area_top + fig_area_h + Inches(0.05)
            caption_h = Inches(0.35)
            legend_top = caption_top + caption_h + Inches(0.05)
            legend_h = SLIDE_H - legend_top - Inches(0.4)
        else:
            area_h = SLIDE_H - area_top - Inches(0.75)
            fig_area_h = area_h

        legend_w = Inches(3.1) if spec.get("side_legend") else 0
        fig_area_w = SLIDE_W - Inches(0.6) - (legend_w + Inches(0.2) if legend_w else 0)
        if resolved:
            if spec.get("layout") == "vertical":
                h_each = int(fig_area_h / len(resolved))
                for i, p in enumerate(resolved):
                    self._picture_fit(slide, p, Inches(0.3), area_top + h_each * i, fig_area_w, h_each - Inches(0.1))
            else:
                w_each = int(fig_area_w / len(resolved))
                for i, p in enumerate(resolved):
                    self._picture_fit(slide, p, Inches(0.3) + w_each * i, area_top, w_each - Inches(0.2), fig_area_h)

        caption = spec.get("caption") or spec.get("note")
        if caption:
            caption_top_val = caption_top if bottom_legend else None
            caption_h_val = caption_h if bottom_legend else None
            self._figure_caption(slide, caption, top=caption_top_val, width=fig_area_w, height=caption_h_val)

        if legend_w:
            legend_left = Inches(0.3) + fig_area_w + Inches(0.2)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, legend_left, area_top, legend_w, area_h)  # pyrefly: ignore [bad-argument-type]
            _fill(box, RGBColor(0xF0, 0xF4, 0xFA))
            box.shadow.inherit = False
            _, tf = _textbox(
                slide,
                legend_left + Inches(0.2),
                area_top + Inches(0.15),
                legend_w - Inches(0.4),
                area_h - Inches(0.3),  # pyrefly: ignore [unsupported-operation]
            )
            p = tf.paragraphs[0]
            p.text = "Shared axis labels"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            for para in spec.get("side_legend_text", "").split("\n\n"):
                pp = tf.add_paragraph()
                pp.text = para
                pp.font.size = Pt(12)
                pp.font.color.rgb = DARK
                pp.space_before = Pt(10)
        elif bottom_legend:
            legend_left = Inches(0.3)
            legend_w = fig_area_w
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, legend_left, legend_top, legend_w, legend_h)  # pyrefly: ignore [bad-argument-type]
            _fill(box, RGBColor(0xF0, 0xF4, 0xFA))
            box.shadow.inherit = False
            _, tf = _textbox(
                slide,
                legend_left + Inches(0.2),
                legend_top + Inches(0.1),  # pyrefly: ignore [unsupported-operation]
                legend_w - Inches(0.4),
                legend_h - Inches(0.2),  # pyrefly: ignore [unsupported-operation]
            )
            p = tf.paragraphs[0]
            p.text = "Shared axis labels"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT

            legend_text = spec.get("side_legend_text") or spec.get("bottom_legend_text", "")
            for para in legend_text.split("\n\n"):
                pp = tf.add_paragraph()
                pp.text = para.replace("\n", " ")
                pp.font.size = Pt(10.5)
                pp.font.color.rgb = DARK
                pp.space_before = Pt(3)
        self._record_script(spec["title"], [caption or ""])

    # ── Slides ──────────────────────────────────────────────────────────────────

    def cover(self) -> None:
        slide = self._new_slide()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)  # pyrefly: ignore [bad-argument-type]
        _fill(bg, DARK)
        title_len = len(self.content["title"])
        title_sz = 26 if title_len > 80 else 32 if title_len > 55 else 36
        band_top = Inches(4.15)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, band_top, SLIDE_W, Inches(0.06))  # pyrefly: ignore [bad-argument-type]
        _fill(band, ACCENT)
        _, tf = _textbox(slide, Inches(0.9), Inches(1.15), SLIDE_W - Inches(1.8), Inches(2.85))
        p = tf.paragraphs[0]
        p.text = self.content["title"]
        p.font.size = Pt(title_sz)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = self.content["subtitle"]
        p2.font.size = Pt(17)
        p2.font.color.rgb = LIGHT_TXT
        p2.space_before = Pt(10)
        _, tf2 = _textbox(slide, Inches(0.9), band_top + Inches(0.25), SLIDE_W - Inches(1.8), Inches(3.0))
        p3 = tf2.paragraphs[0]
        p3.text = self.author
        p3.font.size = Pt(19)
        p3.font.bold = True
        p3.font.color.rgb = WHITE
        if self.coauthors:
            p4 = tf2.add_paragraph()
            p4.text = "with " + ", ".join(self.coauthors)
            p4.font.size = Pt(14)
            p4.font.color.rgb = LIGHT_TXT
            p4.space_before = Pt(6)
        if self.groups:
            p5 = tf2.add_paragraph()
            p5.text = "   ·   ".join(self.groups)
            p5.font.size = Pt(11)
            p5.font.italic = True
            p5.font.color.rgb = RGBColor(0x8A, 0x9B, 0xB0)
            p5.space_before = Pt(10)
        self._record_script(
            self.content["title"],
            [f"Presented by {self.author}." + (f" With {', '.join(self.coauthors)}." if self.coauthors else "")],
        )

    def agenda(self) -> None:
        slide = self._new_slide()
        self._title_bar(slide, "Agenda")
        items = self.content["agenda"]
        half = (len(items) + 1) // 2
        for col, chunk in enumerate([items[:half], items[half:]]):
            _, tf = _textbox(
                slide,
                Inches(0.7) + col * int((SLIDE_W - Inches(1.4)) / 2),
                Inches(1.4),
                int((SLIDE_W - Inches(1.6)) / 2),
                SLIDE_H - Inches(1.8),
            )
            for i, item in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{col * half + i + 1}.  {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                p.space_after = Pt(12)
        self._record_script("Agenda", ["Today's agenda: " + "; ".join(items) + "."])

    def content_slide(self, key: str) -> None:
        spec = self.content["slides"][key]
        slide = self._new_slide()
        self._title_bar(slide, spec["title"])
        fig_path = self.figures_dir / spec["figure"] if spec.get("figure") else None
        has_fig = bool(fig_path) and fig_path.exists()
        if spec.get("figure") and not has_fig:
            print(f"  [WARN] Figure not found: {fig_path}")

        if spec.get("equation"):
            col_w = int(SLIDE_W / 2) - Inches(0.4) if has_fig else SLIDE_W - Inches(1.2)
            eq_left = Inches(0.6)
            bullets_top = self._equation_focus(
                slide, spec["equation"], left=eq_left, width=col_w, size_pt=15 if has_fig else 20
            )
            if spec.get("caption"):
                self._equation_caption(slide, spec["caption"], top=bullets_top, left=eq_left, width=col_w)
                bullets_top += Inches(0.85)
            self._bullets(
                slide, spec["bullets"], left=eq_left, top=bullets_top, width=col_w, size=13 if has_fig else 14
            )
            if has_fig:
                self._picture_fit(
                    slide,
                    fig_path,
                    int(SLIDE_W / 2) + Inches(0.2),
                    Inches(1.2),
                    int(SLIDE_W / 2) - Inches(0.8),
                    SLIDE_H - Inches(1.7),
                )
        elif spec.get("diagram") == "pipeline":
            self._pipeline_diagram(slide)
            bullets_top = Inches(5.6)
            if spec.get("caption"):
                self._figure_caption(slide, spec["caption"], top=bullets_top)
                bullets_top += Inches(0.75)
            self._bullets(slide, spec["bullets"], top=bullets_top, size=14)
        elif spec.get("diagram") == "policy_grid":
            self._policy_grid_diagram(slide)
            bullets_top = Inches(5.85)
            if spec.get("caption"):
                self._figure_caption(slide, spec["caption"], top=bullets_top)
                bullets_top += Inches(0.6)
            self._bullets(slide, spec["bullets"], top=bullets_top, size=13)
        elif spec.get("diagram") == "doe_tree":
            self._doe_tree_diagram(slide)
            bullets_top = Inches(5.85)
            if spec.get("caption"):
                self._figure_caption(slide, spec["caption"], top=bullets_top)
                bullets_top += Inches(0.6)
            self._bullets(slide, spec["bullets"], top=bullets_top, size=13)
        elif has_fig:
            self._picture_fit(
                slide, fig_path, int(SLIDE_W / 2), Inches(1.2), int(SLIDE_W / 2) - Inches(0.4), SLIDE_H - Inches(1.7)
            )
            if spec.get("caption"):
                self._figure_caption(slide, spec["caption"])
            self._bullets(
                slide, spec["bullets"], left=Inches(0.5), top=Inches(1.4), width=int(SLIDE_W / 2) - Inches(0.7), size=13
            )
        else:
            self._bullets(slide, spec["bullets"])
        self._record_script(spec["title"], [spec.get("caption", "")] + list(spec["bullets"]))

    def _load_horizon(self, spec: dict) -> tuple | None:
        csv_path = Path(spec["csv"])
        if not csv_path.exists():
            return None
        df = filter_data(load_horizon_csv(csv_path), self._results_config)
        if df.empty:
            return None
        theme = {"name": "light", "fg": "#1F2D3D"}
        ctx = build_context(df, self._results_config, theme, spec["days"])
        return df, aggregate(df), ctx

    def _results_table_slide(self) -> None:
        """Full results table slide, for the horizon(s) chosen via --results-table."""
        if self.results_table == "none":
            return
        self._results_config = load_json("simulation_analysis_config.json")
        horizon_specs = sorted(self._results_config["horizons"], key=lambda h: h["days"])

        if self.results_table == "all":
            row_keys, col_keys, cells = [], [], {}
            used_days = []
            for spec in horizon_specs:
                loaded = self._load_horizon(spec)
                if loaded is None:
                    continue
                _, dfm, ctx = loaded
                rk, ck, cs = build_full_results_matrix(dfm, ctx, horizon_label=f"{spec['days']}d")
                row_keys += [r for r in rk if r not in row_keys]
                col_keys += ck
                cells.update(cs)
                used_days.append(spec["days"])
            row_keys.sort()
            if not row_keys:
                print("  [WARN] No data available for the full results table — skipping slide")
                return
            title = f"Results — Full Table (All Horizons: {', '.join(f'{d}d' for d in used_days)})"
            multi = True
        else:
            days = int(self.results_table.rstrip("d"))
            spec = next((s for s in horizon_specs if s["days"] == days), None)
            loaded = self._load_horizon(spec) if spec else None
            if loaded is None:
                print(f"  [WARN] No data for the {self.results_table} horizon — skipping full results table slide")
                return
            _, dfm, ctx = loaded
            row_keys, col_keys, cells = build_full_results_matrix(dfm, ctx)
            title = f"Results — Full Table ({days}-Day Horizon)"
            multi = False

        level_phrases = {
            "horizon": "horizon",
            "strategy": "mandatory selection strategy",
            "constructor": "route constructor",
            "improver": "route improver",
        }
        level_names = (["horizon"] if multi else []) + ["strategy", "constructor", "improver"]
        col_desc = " × ".join(level_phrases[n] for n in level_names)

        slide = self._new_slide()
        self._title_bar(slide, title)
        area_top, area_bottom = Inches(1.15), SLIDE_H - Inches(0.95)
        split = self.results_table_split

        if split != "none" and split in level_names:
            level_idx = level_names.index(split)
            partition_values: list = []
            for ck in col_keys:
                if ck[level_idx] not in partition_values:
                    partition_values.append(ck[level_idx])
            part_h = int((area_bottom - area_top) / len(partition_values))
            for pi, val in enumerate(partition_values):
                sub_full = [ck for ck in col_keys if ck[level_idx] == val]
                sub_display = [ck[:level_idx] + ck[level_idx + 1 :] for ck in sub_full]
                part_top = area_top + pi * part_h
                # Embed the partition label directly in the table image — no separate textbox
                img_path = render_hier_table_image(
                    row_keys,
                    sub_display,
                    cells,
                    ["N", "Region", "Dist"],
                    self._tmp / f"full_results_table_{split}_{pi}.png",
                    col_lookup_keys=sub_full,
                    partition_label=f"{split.capitalize()}: {val}",
                )
                self._picture_fit(
                    slide,
                    img_path,
                    Inches(0.3),
                    part_top,
                    SLIDE_W - Inches(0.6),
                    part_h,
                )
            remaining_desc = " × ".join(level_phrases[n] for n in level_names if n != split)
            col_desc = (
                f"partitioned by {level_phrases[split]} ({', '.join(str(v) for v in partition_values)}), "
                f"one partial table per value; each partial table then grouped by {remaining_desc}"
            )
        else:
            img_path = render_hier_table_image(
                row_keys, col_keys, cells, ["N", "Region", "Dist"], self._tmp / "full_results_table.png"
            )
            self._picture_fit(slide, img_path, Inches(0.3), area_top, SLIDE_W - Inches(0.6), area_bottom - area_top)

        self._table_caption(
            slide,
            f"Full results table — rows: graph size × region × data distribution; columns: {col_desc}. "
            "In every cell, the top value is mean±std overflows (lower is better) and the bottom value is "
            "mean±std kg/km (higher is better); the best value per row is shown in bold green.",
        )
        self._record_script(title, [f"Full results table, columns grouped by {col_desc}."])

    def qa(self) -> None:
        spec = self.content["slides"]["qa"]
        slide = self._new_slide()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)  # pyrefly: ignore [bad-argument-type]
        _fill(bg, DARK)
        fig_path = self.figures_dir / spec["figure"] if spec.get("figure") else None
        has_fig = bool(fig_path) and fig_path.exists()
        text_left = Inches(0.9)
        text_w = (int(SLIDE_W / 2) - Inches(1.1)) if has_fig else (SLIDE_W - Inches(1.8))
        _, tf = _textbox(slide, text_left, Inches(2.8), text_w, Inches(2.4))
        p = tf.paragraphs[0]
        p.text = spec["title"]
        p.font.size = Pt(30 if has_fig else 36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER if not has_fig else PP_ALIGN.LEFT
        for line in spec["bullets"]:
            pp = tf.add_paragraph()
            pp.text = line
            pp.font.size = Pt(16)
            pp.font.color.rgb = LIGHT_TXT
            pp.alignment = PP_ALIGN.CENTER if not has_fig else PP_ALIGN.LEFT
        if has_fig:
            self._picture_fit(
                slide,
                fig_path,
                int(SLIDE_W / 2) + Inches(0.2),
                Inches(0.6),
                int(SLIDE_W / 2) - Inches(0.8),
                SLIDE_H - Inches(1.2),
            )
        self._record_script(spec["title"], list(spec["bullets"]))

    def build(self) -> PresentationClass:
        self.cover()  # 1
        self.agenda()  # 2
        self.content_slide("vrpp")  # 3
        self.content_slide("objective")  # 4
        self.content_slide("simulator")  # 5
        self.content_slide("strategies")  # 6
        self.content_slide("policy_overview")  # 7
        self.content_slide("exact")  # 8
        self.content_slide("metaheuristics")  # 9
        self.content_slide("improvers")  # 10
        self.content_slide("design_of_experiments")  # 11
        self._figure_slide(self.content["figure_slides"]["pareto"])  # 12
        self._figure_slide(self.content["figure_slides"]["kpi"])  # 13
        self._figure_slide(self.content["figure_slides"]["strategy_bubble"])  # 14
        self._figure_slide(self.content["figure_slides"]["scenario_heatmaps"])  # 15 (30d)
        self._figure_slide(self.content["figure_slides"]["heatmaps"])  # 16 (90d)
        self._figure_slide(self.content["figure_slides"]["improver_bubble"])  # 17
        self._results_table_slide()  # 18 (full results table, if requested)
        self.content_slide("conclusion")  # 19
        self.qa()  # 20
        return self.prs


def generate_qa_route_image(out_path: Path) -> Path:
    """Generate a VRP with Profits routing illustration with multiple routes and unvisited nodes."""
    # Set seed for reproducibility
    np.random.seed(42)

    # Generate nodes
    n_nodes = 45
    coords = np.random.rand(n_nodes, 2)
    depot = np.array([0.5, 0.5])

    # Calculate angles and distances from depot
    diffs = coords - depot
    angles = np.arctan2(diffs[:, 1], diffs[:, 0])
    distances = np.linalg.norm(diffs, axis=1)

    route1_idx = []
    route2_idx = []
    route3_idx = []
    unvisited_idx = []

    for i in range(n_nodes):
        # Leave about 30% of the nodes unvisited to show the profit-based node selection
        if i % 3 == 0 and distances[i] > 0.15:
            unvisited_idx.append(i)
            continue

        angle = angles[i]
        if -np.pi <= angle < -np.pi / 3:
            route1_idx.append(i)
        elif -np.pi / 3 <= angle < np.pi / 3:
            route2_idx.append(i)
        else:
            route3_idx.append(i)

    # Sort nodes in each route by polar angle to make smooth tours
    def sort_route(indices):
        if not indices:
            return []
        sub_angles = angles[indices]
        sorted_indices = [indices[idx] for idx in np.argsort(sub_angles)]
        return sorted_indices

    r1 = sort_route(route1_idx)
    r2 = sort_route(route2_idx)
    r3 = sort_route(route3_idx)

    fig, ax = plt.subplots(figsize=(8, 8), dpi=200)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    # Draw faint gray edges between nearby nodes to represent the road network
    for i in range(n_nodes):
        for j in range(i + 1, n_nodes):
            dist = np.linalg.norm(coords[i] - coords[j])
            if dist < 0.22:
                ax.plot([coords[i, 0], coords[j, 0]], [coords[i, 1], coords[j, 1]], color="#E2E8F0", lw=0.6, zorder=1)

    # Plot unvisited nodes (selection part of VRPP)
    ax.scatter(
        coords[unvisited_idx, 0],
        coords[unvisited_idx, 1],
        color="#CBD5E1",
        edgecolor="#64748B",
        s=100,
        label="Unvisited Bins (No Profit)",
        zorder=3,
    )

    # Premium colors for the 3 routes (different vehicles)
    colors = ["#10B981", "#6366F1", "#F59E0B"]
    route_labels = ["Vehicle Route 1", "Vehicle Route 2", "Vehicle Route 3"]

    # Plot each route
    def plot_tour(route_nodes, color, label):
        if not route_nodes:
            return
        tour_coords = [depot] + [coords[idx] for idx in route_nodes] + [depot]
        tour_coords = np.array(tour_coords)

        ax.plot(tour_coords[:, 0], tour_coords[:, 1], color=color, lw=3, label=label, zorder=2)
        ax.scatter(tour_coords[1:-1, 0], tour_coords[1:-1, 1], color=color, edgecolor="black", s=120, zorder=4)

    plot_tour(r1, colors[0], route_labels[0])
    plot_tour(r2, colors[1], route_labels[1])
    plot_tour(r3, colors[2], route_labels[2])

    # Plot depot clearly marked
    ax.scatter(
        depot[0],
        depot[1],
        color="#EF4444",
        marker="s",
        s=220,
        edgecolor="black",
        linewidth=2,
        label="Depot / Warehouse",
        zorder=5,
    )

    ax.set_xlim(-0.05, 1.05)
    ax.set_ylim(-0.05, 1.05)
    ax.axis("off")

    ax.legend(
        loc="lower center",
        bbox_to_anchor=(0.5, -0.08),
        ncol=2,
        frameon=True,
        facecolor="white",
        edgecolor="#E2E8F0",
        fontsize=10,
    )

    plt.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument(
        "--figures-dir",
        default="public/figures/simulation/30d",
        help="Directory with the simulation analysis figures to embed",
    )
    p.add_argument("--out", default="assets/windows/wsmart_route_results.pptx", help="Destination .pptx path")
    p.add_argument("--author", default=None, help="Presenting author shown on the cover slide")
    p.add_argument(
        "--coauthors", default=None, help="Optional semicolon-separated co-authors (shown with less emphasis)"
    )
    p.add_argument("--groups", default=None, help="Optional semicolon-separated research groups of the authors")
    p.add_argument(
        "--results-table",
        default="30d",
        choices=["30d", "90d", "all", "none"],
        help="Full results table slide: a single horizon, 'all' horizons "
        "(horizon added to the column hierarchy), or 'none' to omit the slide",
    )
    p.add_argument(
        "--results-table-split",
        default="none",
        choices=["none", "strategy", "constructor", "improver"],
        help="Partition the results table by this column-hierarchy level: it becomes the "
        "outermost grouping and each value gets its own stacked partial table on the slide "
        "(e.g. 'improver' -> one partial table for CLS, one for FTSP)",
    )
    p.add_argument(
        "--speaker-script",
        action="store_true",
        help="Also generate a per-slide speaker script .docx alongside the .pptx",
    )
    p.add_argument(
        "--speaker-script-out",
        default=None,
        help="Destination .docx path (default: same name as --out, .docx extension)",
    )
    return p.parse_args()


def main() -> None:
    args = parse_args()
    content = load_json("presentation_content.json")
    figures_dir = Path(args.figures_dir)
    if not figures_dir.is_dir():
        raise SystemExit(f"Figures dir not found: {figures_dir} — run gen_simulation_analysis.py first")

    # Generate the QA route illustration dynamically
    generate_qa_route_image(figures_dir / "qa_route_illustration.png")

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    coauthors = [c.strip() for c in args.coauthors.split(";")] if args.coauthors else None
    groups = [g.strip() for g in args.groups.split(";")] if args.groups else None
    builder = DeckBuilder(
        content, figures_dir, args.author, coauthors, groups, args.results_table, args.results_table_split
    )
    prs = builder.build()
    prs.save(str(out))
    print(f"Written: {out} ({len(prs.slides._sldIdLst)} slides)")

    if args.speaker_script:
        script_out = Path(args.speaker_script_out) if args.speaker_script_out else out.with_suffix(".docx")
        gen_speaker_script(content["title"], builder.author, builder.slide_scripts, script_out)
        print(f"Written: {script_out} ({len(builder.slide_scripts)} slide scripts)")


if __name__ == "__main__":
    main()
